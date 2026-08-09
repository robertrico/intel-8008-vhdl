#!/usr/bin/env python3
"""Build the FPGA-focused talk: docs/fpga-presentation/fpga_talk.pptx

Sibling of docs/presentation/deck_build.py (the 8008-focused deck). Same
design system, different center of gravity: this deck is about FPGAs and
HDL, and uses the 8008 only as the payload that makes the FPGA content
concrete.

Theme: Lattice charcoal + yellow is the DEFAULT zone here (we live in the
FPGA world the whole talk). Intel blue appears only on the two "what is the
payload" slides, as a guest.

Run:  python3 deck_build.py       (needs python-pptx + Pillow)
"""
import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = '/Users/hambook/Development/intel-8008-vhdl'
ASSETS = REPO + '/docs/fpga-presentation/assets'
OUT = REPO + '/docs/fpga-presentation/fpga_talk.pptx'

INTEL = RGBColor(0x00, 0x68, 0xB5)
INTEL_DARK = RGBColor(0x00, 0x42, 0x77)
ENERGY = RGBColor(0x00, 0xC7, 0xFD)
YELLOW = RGBColor(0xFF, 0xC2, 0x20)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
CHARCOAL2 = RGBColor(0x3B, 0x3F, 0x44)
LIGHTBG = RGBColor(0xF7, 0xF8, 0xFA)
SOFTWHITE = RGBColor(0xF2, 0xF2, 0xF2)
CODEBG = RGBColor(0x2A, 0x2D, 0x31)
SKY = RGBColor(0x1F, 0x9D, 0xD8)
MIDGREY = RGBColor(0x8D, 0x8D, 0x8D)
CODE_COMMENT = RGBColor(0x6F, 0xC2, 0xEA)
PALE_BLUE = RGBColor(0xCD, 0xEB, 0xFA)
PALE_YELLOW = RGBColor(0xFF, 0xF3, 0xD3)
GREYFILL = RGBColor(0xE9, 0xEC, 0xF0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xB3, 0x33, 0x2E)

FONT = 'Aptos'
MONO = 'Aptos Mono'

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

ZONE = 'lattice'   # this deck lives in FPGA-land; 'intel' is the guest zone


# ------------------------------------------------------------------ helpers

def z_dark_bg():
    return INTEL_DARK if ZONE == 'intel' else CHARCOAL


def z_primary():
    return INTEL if ZONE == 'intel' else CHARCOAL


def z_accent_dark():
    return ENERGY if ZONE == 'intel' else YELLOW


def z_accent_light():
    return INTEL if ZONE == 'intel' else YELLOW


def z_code_inline():
    return INTEL if ZONE == 'intel' else SKY


def have(path):
    return os.path.exists(path)


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = z_dark_bg() if dark else LIGHTBG
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def style_run(r, dark, size, mono=False, bold=False, color=None):
    r.font.name = MONO if mono else FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color if color else (SOFTWHITE if dark else CHARCOAL)


def rich(p, text, dark, size, base_color=None):
    """Parse **bold** and `mono` runs into paragraph p."""
    for tok in re.split(r'(\*\*.*?\*\*|`.*?`)', text):
        if not tok:
            continue
        if tok.startswith('**') and tok.endswith('**'):
            r = p.add_run(); r.text = tok[2:-2]
            style_run(r, dark, size, bold=True, color=base_color)
        elif tok.startswith('`') and tok.endswith('`'):
            r = p.add_run(); r.text = tok[1:-1]
            style_run(r, dark, size, mono=True,
                      color=z_accent_dark() if dark else z_code_inline())
        else:
            r = p.add_run(); r.text = tok
            style_run(r, dark, size, color=base_color)


def bullets(s, x, y, w, h, lines, dark=False, size=18, space=8, anchor_top=True):
    tb, tf = box(s, x, y, w, h)
    if not anchor_top:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space)
        indent = line.startswith('  ')
        txt = line.strip()
        if txt.startswith('- '):
            txt = ('–  ' if indent else '•  ') + txt[2:]
            if indent:
                p.level = 1
        rich(p, txt, dark, size - (2 if indent else 0))
    return tb


def kicker(s, text, dark=False):
    tb, tf = box(s, 0.55, 0.32, 8.0, 0.35)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text.upper()
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
    if dark:
        r.font.color.rgb = z_accent_dark()
    else:
        r.font.color.rgb = INTEL if ZONE == 'intel' else CHARCOAL2


def title(s, text, dark=False, size=32, y=0.62, w=12.2):
    tb, tf = box(s, 0.55, y, w, 1.0)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = SOFTWHITE if dark else z_primary()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y + 0.78),
                             Inches(1.6), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = z_accent_dark() if dark else z_accent_light()
    bar.line.fill.background()
    return tb


def rect(s, x, y, w, h, fill, line_color=None, line_w=None,
         shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w or 1)
    sh.shadow.inherit = False
    return sh


def shape_text(sh, text, size=14, bold=False, color=None, mono=False,
               align=PP_ALIGN.CENTER):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = MONO if mono else FONT
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color if color else SOFTWHITE


def placeholder(s, x, y, w, h, label, dark=False):
    sh = rect(s, x, y, w, h, CHARCOAL2 if dark else RGBColor(0xEA, 0xED, 0xF1),
              line_color=MIDGREY, line_w=1.25)
    from pptx.oxml.ns import qn
    ln = sh.line._get_or_add_ln()
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(d)
    shape_text(sh, 'PLACEHOLDER\n' + label, size=13, bold=False,
               color=SOFTWHITE if dark else RGBColor(0x6A, 0x6F, 0x76))
    return sh


def code_panel(s, x, y, w, h, code, size=14, bg=None, comment='--'):
    """Monospace panel. `comment` = the line-comment token to colorize,
    or None for a literal panel (shell commands, reports)."""
    sh = rect(s, x, y, w, h, bg or CODEBG)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.12)
    for i, ln in enumerate(code.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        if comment and comment in ln:
            main, com = ln.split(comment, 1)
            com = comment + com
        else:
            main, com = ln, None
        if main:
            r = p.add_run(); r.text = main
            r.font.name = MONO; r.font.size = Pt(size)
            r.font.color.rgb = SOFTWHITE
        if com is not None:
            r = p.add_run(); r.text = com
            r.font.name = MONO; r.font.size = Pt(size)
            r.font.color.rgb = CODE_COMMENT
        if not main and com is None:
            r = p.add_run(); r.text = ''
            r.font.name = MONO; r.font.size = Pt(size)
    return sh


def caption(s, x, y, w, text, size=12, color=None):
    tb, tf = box(s, x, y, w, 0.4)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.italic = True
    r.font.color.rgb = color or MIDGREY
    return tb


def pic(s, path, x, y, w=None, h=None, border=True):
    kw = {}
    if w: kw['width'] = Inches(w)
    if h: kw['height'] = Inches(h)
    p = s.shapes.add_picture(path, Inches(x), Inches(y), **kw)
    if border:
        p.line.color.rgb = MIDGREY
        p.line.width = Pt(0.75)
    p.shadow.inherit = False
    return p


def img_size(path):
    return Image.open(path).size


def fit(path, max_w, max_h):
    w, h = img_size(path)
    scale = min(max_w / w, max_h / h)
    return w * scale, h * scale


def logo(s, path, h=0.4, x=None, y=0.35):
    iw, ih = img_size(path)
    w = h * iw / ih
    if x is None:
        x = 13.333 - w - 0.55
    return pic(s, path, x, y, h=h, border=False)


def arrow(s, x1, y1, x2, y2, color=None, wpt=2.25):
    conn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color if color else z_primary()
    conn.line.width = Pt(wpt)
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    head = ln.makeelement(qn('a:tailEnd'), {'type': 'arrow'})
    ln.append(head)
    conn.shadow.inherit = False
    return conn


def line(s, x1, y1, x2, y2, color=None, wpt=1.25):
    conn = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color or MIDGREY
    conn.line.width = Pt(wpt)
    conn.shadow.inherit = False
    return conn


def table(s, x, y, w, h, rows, col_widths, header=True, size=16):
    shp = s.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y),
                             Inches(w), Inches(h))
    t = shp.table
    hdr_color = INTEL if ZONE == 'intel' else CHARCOAL2
    alt = RGBColor(0xED, 0xF1, 0xF6) if ZONE == 'intel' else RGBColor(0xEF, 0xEF, 0xEF)
    for ci, cw in enumerate(col_widths):
        t.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            if header and ri == 0:
                cell.fill.fore_color.rgb = hdr_color
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 else alt
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            rich(p, str(val), False, size,
                 base_color=RGBColor(0xFF, 0xFF, 0xFF) if (header and ri == 0) else CHARCOAL)
            if header and ri == 0:
                for r in p.runs:
                    r.font.bold = True
    return t


def stat_tiles(s, x, y, items, tile_w=2.6, tile_h=1.5, gap=0.25, dark=False,
               num_size=30, lbl_size=12):
    """Row of big-number tiles: items = [(number, label), ...]"""
    cx = x
    for num, lbl in items:
        sh = rect(s, cx, y, tile_w, tile_h,
                  CHARCOAL2 if dark else RGBColor(0xFF, 0xFF, 0xFF),
                  line_color=None if dark else RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.name = FONT; r.font.size = Pt(num_size); r.font.bold = True
        r.font.color.rgb = YELLOW if dark else CHARCOAL
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        r2 = p2.add_run(); r2.text = lbl
        r2.font.name = FONT; r2.font.size = Pt(lbl_size)
        r2.font.color.rgb = SOFTWHITE if dark else RGBColor(0x5A, 0x60, 0x68)
        if not dark:
            rect(s, cx, y, 0.06, tile_h, YELLOW)
        cx += tile_w + gap


def chain(s, x, y, labels, box_w=1.9, box_h=1.05, gap=0.52, dark=False,
          fills=None, size=13):
    """Horizontal flow: boxes with arrows between."""
    cx = x
    shapes = []
    for i, lbl in enumerate(labels):
        fill = (fills[i] if fills else (CHARCOAL2 if dark else RGBColor(0xFF, 0xFF, 0xFF)))
        sh = rect(s, cx, y, box_w, box_h, fill,
                  line_color=None if dark else RGBColor(0xC9, 0xCE, 0xD6), line_w=1.25,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(sh, lbl, size=size, bold=True,
                   color=SOFTWHITE if dark else CHARCOAL)
        shapes.append(sh)
        if i < len(labels) - 1:
            arrow(s, cx + box_w + 0.06, y + box_h / 2, cx + box_w + gap - 0.06,
                  y + box_h / 2, color=YELLOW if dark else CHARCOAL2, wpt=2.0)
        cx += box_w + gap
    return shapes


# ============================================================
# FRAME
# ============================================================

# ---- Slide 1: Title
s = slide(dark=True)
tb, tf = box(s, 0.75, 1.85, 7.6, 3.8)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'Nothing → CPU'
r.font.name = MONO; r.font.size = Pt(46); r.font.bold = True
r.font.color.rgb = YELLOW
p = tf.add_paragraph(); p.space_before = Pt(10)
r = p.add_run(); r.text = 'What an FPGA actually is,\nand what your HDL becomes inside it'
r.font.name = FONT; r.font.size = Pt(30); r.font.bold = True
r.font.color.rgb = SOFTWHITE
p = tf.add_paragraph(); p.space_before = Pt(20)
r = p.add_run(); r.text = ('Told through a working Intel 8008 in VHDL on a Lattice ECP5 —\n'
                          'built with a fully open-source toolchain')
r.font.name = FONT; r.font.size = Pt(17)
r.font.color.rgb = RGBColor(0xC9, 0xCE, 0xD6)
p = tf.add_paragraph(); p.space_before = Pt(24)
r = p.add_run(); r.text = 'Robert Rico'
r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
r.font.color.rgb = SOFTWHITE
rect(s, 0.8, 1.78, 1.6, 0.045, YELLOW)
logo(s, ASSETS + '/lattice_logo_ondark.png', h=0.42, x=8.95, y=1.95)
if have(REPO + '/docs/images/tiny_os_scelbal.png'):
    tw, th = fit(REPO + '/docs/images/tiny_os_scelbal.png', 3.85, 2.8)
    pic(s, REPO + '/docs/images/tiny_os_scelbal.png', 8.95, 3.0, w=tw)
    caption(s, 8.95, 3.0 + th + 0.06, 3.85,
            'The finish line: it boots into BASIC.', color=RGBColor(0x9A, 0xA0, 0xA8))
notes(s, 'Open by naming the real subject: this is an FPGA and HDL talk. The 8008 is the '
         'payload — it is what makes the FPGA content concrete instead of abstract. '
         'One sentence of setup: "I built a 1972 microprocessor out of an empty chip, and '
         'along the way I had to learn what that chip actually is. That is the talk."')

# ---- Slide 2: Roadmap
s = slide()
kicker(s, 'Roadmap')
title(s, 'Five questions, in order')
cards = [
    ('1', 'What IS an\nFPGA?', 'LUTs, flip-flops,\nrouting, hard blocks,\nbitstreams. No\nhand-waving.'),
    ('2', 'What does HDL\nactually say?', 'VHDL vs Verilog, the\nthree things you write,\nand what each becomes\nin silicon.'),
    ('3', 'What does the\ntoolchain do?', 'GHDL → Yosys →\nnextpnr → bitstream.\nSynthesis, place, route,\ntiming closure.'),
    ('4', 'Where does it\nbite you?', 'Tri-states, inference,\nclock domains, and the\nday simulation lied\nto me.'),
    ('5', 'Can you learn\nthis with an AI?', 'Yes. I did. And the two\nattempts I threw away\nshow exactly how\nnot to.'),
]
cx = 0.6
for num, head, body in cards:
    sh = rect(s, cx, 2.0, 2.35, 4.4, RGBColor(0xFF, 0xFF, 0xFF),
              line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
    rect(s, cx, 2.0, 2.35, 0.09, YELLOW)
    tf2 = sh.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.2); tf2.margin_right = Inches(0.14)
    tf2.margin_top = Inches(0.26)
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = num
    r.font.name = MONO; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = RGBColor(0xC9, 0xCE, 0xD6)
    p = tf2.add_paragraph(); p.space_before = Pt(6); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = head
    r.font.name = FONT; r.font.size = Pt(16.5); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    p = tf2.add_paragraph(); p.space_before = Pt(10); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = body
    r.font.name = FONT; r.font.size = Pt(12.5)
    r.font.color.rgb = RGBColor(0x5A, 0x60, 0x68)
    cx += 2.48   # 5 cards -> right edge 12.87
bullets(s, 0.6, 6.65, 12.2, 0.7, [
    'Then: a live 1972 machine, running as a circuit inside a 2013 chip.',
], size=17)
notes(s, 'Set the contract. Software people get 1–2; hardware people get 3–4; question 5 '
         'is the one that starts arguments, so flag it early and honestly rather than '
         'springing it. Say out loud: "I had never used an FPGA and could not read VHDL '
         'when I started. I am going to show you what worked and, more usefully, what '
         'did not."')

# ============================================================
# PART 1 — WHAT AN FPGA IS
# ============================================================

# ---- Slide 3: Three ways to build a digital system
s = slide()
kicker(s, 'Part 1 — What an FPGA is')
title(s, 'Three ways to build a digital system')
table(s, 0.6, 1.95, 12.1, 3.2, [
    ['', 'CPU + software', 'FPGA', 'ASIC'],
    ['What it is', 'Fixed circuit, running instructions',
     '**Blank circuit you configure**', 'Custom-etched silicon'],
    ['Change it by', 'Recompiling (seconds)', 'Re-synthesizing (minutes)',
     'A new mask set ($$$, months)'],
    ['Concurrency', 'A few cores', '**Everything, at once**', 'Everything, at once'],
    ['Cost to start', '$0', '$50 – $5,000 board', '$1M+'],
    ['Good for', 'Anything sequential', 'Parallel, timing-exact, custom I/O',
     'Volume, power, absolute speed'],
], [1.9, 3.1, 3.6, 3.5], size=14)
sh = rect(s, 0.6, 5.5, 12.1, 1.15, CHARCOAL)
tf = sh.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'An FPGA gives you **the flexibility of software with the behavior of '
        'hardware.** You do not write a program. You wire up a machine.', True, 18)
notes(s, 'The single most useful framing for a software audience. An FPGA is not a fast '
         'CPU and it is not a simulator — it is a bag of gates you get to arrange, and '
         'then rearrange tomorrow. That "rearrange tomorrow" is the whole value '
         'proposition; ASIC people only get one shot.')

# ---- Slide 4: Inside the chip — the fabric
s = slide()
kicker(s, 'Part 1 — What an FPGA is')
title(s, 'Inside: a grid of tiny identical tiles')
# fabric grid
gx, gy, cell, gap = 0.65, 2.05, 0.52, 0.09
cols, rows = 9, 6
for rr in range(rows):
    for cc in range(cols):
        is_bram = cc in (3, 7)
        fill = PALE_YELLOW if is_bram else RGBColor(0xFF, 0xFF, 0xFF)
        edge = YELLOW if is_bram else RGBColor(0xCF, 0xD4, 0xDB)
        rect(s, gx + cc * (cell + gap), gy + rr * (cell + gap), cell, cell,
             fill, line_color=edge, line_w=1)
# highlight one logic tile
hx = gx + 1 * (cell + gap)
hy = gy + 2 * (cell + gap)
rect(s, hx - 0.04, hy - 0.04, cell + 0.08, cell + 0.08, None,
     line_color=CHARCOAL, line_w=2.5)
caption(s, 0.65, gy + rows * (cell + gap) + 0.08, 5.6,
        'ECP5 fabric: logic tiles (white), block-RAM columns (yellow), '
        'I/O ring around the edge')
# zoom panel: one slice
zx = 6.6
rect(s, zx, 1.98, 6.1, 4.05, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1.25)
tb, tf = box(s, zx + 0.25, 2.08, 5.6, 0.5)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'One tile, unpacked'
r.font.name = FONT; r.font.size = Pt(17); r.font.bold = True
r.font.color.rgb = CHARCOAL
lut = rect(s, zx + 0.45, 2.85, 1.85, 1.35, CHARCOAL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(lut, 'LUT4\nany function\nof 4 inputs', size=13, bold=True)
ff = rect(s, zx + 3.05, 2.85, 1.5, 1.35, CHARCOAL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(ff, 'FF\nremembers\none bit', size=13, bold=True)
arrow(s, zx + 2.32, 3.52, zx + 3.02, 3.52, color=CHARCOAL2, wpt=2)
arrow(s, zx + 4.57, 3.52, zx + 5.35, 3.52, color=CHARCOAL2, wpt=2)
for i in range(4):
    arrow(s, zx + 0.05, 2.98 + i * 0.34, zx + 0.42, 2.98 + i * 0.34,
          color=MIDGREY, wpt=1.5)
cc_ = rect(s, zx + 0.45, 4.45, 4.1, 0.65, PALE_YELLOW,
           line_color=YELLOW, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(cc_, 'dedicated carry chain (CCU2C) — fast adders without burning LUTs',
           size=12, color=CHARCOAL)
bullets(s, zx + 0.45, 5.28, 5.4, 0.6, [
    'Multiply that tile by **43,848**. Add wires between all of them.',
], size=15)
bullets(s, 0.65, 6.35, 12.1, 1.0, [
    'The wires matter as much as the logic: most of an FPGA’s die area — and most of your '
    'timing delay — is **programmable routing**, not gates.',
], size=17)
notes(s, 'Land two ideas: (1) the chip is homogeneous — the same tile, tens of thousands '
         'of times; (2) the interesting part is the routing between them. When a design '
         'fails timing it is almost always routing delay, not logic delay. Point at the '
         'highlighted tile and say: next slide, we open that LUT.')

# ---- Slide 5: The LUT
s = slide()
kicker(s, 'Part 1 — What an FPGA is')
title(s, 'The LUT: a truth table you can rewrite')
bullets(s, 0.6, 1.9, 5.5, 1.9, [
    'A 4-input LUT is **16 bits of memory** plus a 16-to-1 multiplexer.',
    'The four inputs are the address. The stored bit is the answer.',
], size=17, space=12)
code_panel(s, 0.6, 3.55, 5.5, 1.0,
           'q <= (a and b) or (c and not d);', size=15)
bullets(s, 0.6, 4.75, 5.5, 2.2, [
    'No AND gate is built. No OR gate is built.',
    'The synthesizer **evaluates your expression for all 16 input combinations** and '
    'stores the answers.',
], size=16, space=10)
# LUT diagram
lx, ly = 6.7, 1.95
rect(s, lx, ly, 6.0, 4.6, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1.25)
bits = [0, 0, 0, 1, 1, 1, 1, 1, 0, 0, 0, 1, 0, 0, 0, 1]
cw, chh, cg = 0.52, 0.42, 0.08
tbx, tby = lx + 1.55, ly + 0.75
tb, tf = box(s, lx + 0.25, ly + 0.18, 5.5, 0.45)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'The 16 configuration bits your bitstream loads'
r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True
r.font.color.rgb = CHARCOAL
for i, b in enumerate(bits):
    rr, cc = divmod(i, 4)
    sh = rect(s, tbx + cc * (cw + cg), tby + rr * (chh + cg), cw, chh,
              YELLOW if b else RGBColor(0xEC, 0xEF, 0xF3),
              line_color=RGBColor(0xC9, 0xCE, 0xD6), line_w=1)
    shape_text(sh, str(b), size=14, bold=True, color=CHARCOAL, mono=True)
    if cc == 0:
        caption(s, tbx - 1.28, tby + rr * (chh + cg) + 0.04, 1.2,
                f'dcba = {i:04b}..{i + 3:04b}', size=10)
outp = rect(s, tbx + 0.9, tby + 4 * (chh + cg) + 0.35, 2.1, 0.6, CHARCOAL2,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(outp, 'q', size=16, bold=True, mono=True)
arrow(s, tbx + 1.95, tby + 4 * (chh + cg) + 0.05, tbx + 1.95,
      tby + 4 * (chh + cg) + 0.32, color=CHARCOAL2, wpt=2)
caption(s, lx + 0.25, ly + 4.05, 5.5,
        'Change the function → change 16 bits. Same silicon.', size=13)
notes(s, 'This is the "oh!" slide for anyone who has never seen inside an FPGA. Walk one '
         'row: a=1,b=1 → the answer is 1 regardless of c and d, so those rows are yellow. '
         'Then the punchline: reconfiguring an FPGA is literally rewriting these tables. '
         'The ECP5 has ~44,000 of them, and a real LUT4 also has a fast path to its '
         'neighbour so wide functions chain cheaply.')

# ---- Slide 6: The hard blocks
s = slide()
kicker(s, 'Part 1 — What an FPGA is')
title(s, 'LUTs alone would be wasteful — so: hard blocks')
bullets(s, 0.6, 1.85, 12.2, 0.8, [
    'Some things are so common that building them from LUTs is silly. Vendors etch them '
    'in as fixed silicon, scattered through the fabric:',
], size=17)
table(s, 0.6, 2.8, 12.1, 2.9, [
    ['Block', 'ECP5-45F has', 'What it is', 'This project uses'],
    ['**LUT4 / FF**', '43,848', 'The general-purpose fabric', '1,770 LUTs, 591 FFs  (4%)'],
    ['**DP16KD**', '108', '18 kb dual-port block RAM', '**8** — the 8008’s ROM and RAM'],
    ['**EHXPLLL**', '4', 'PLL: makes new clocks from one clock', '**1** — 100 MHz → 25 MHz'],
    ['**MULT18X18D**', '72', 'Hardware multiplier / DSP', '0 — a 1972 CPU has no multiply'],
    ['**TRELLIS_IO**', '245', 'Configurable I/O pin buffers', '63 — UART, LEDs, switches'],
], [2.3, 1.9, 4.3, 3.6], size=14)
sh = rect(s, 0.6, 6.05, 12.1, 1.0, PALE_YELLOW, line_color=YELLOW, line_w=1.5)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'You do not instantiate these by name. You **write a pattern the tool recognizes**, '
        'and it swaps in the hard block. That is called *inference* — and we come back to '
        'it, because getting it wrong is expensive.', False, 16)
notes(s, 'Real numbers from this project’s nextpnr report — utilization.txt in '
         'projects/b8008_basic/reports/. The DP16KD row is the one to dwell on: 8 block '
         'RAMs hold the entire 1970s software library. The PLL row sets up the clocking '
         'slide later.')

# ---- Slide 7: The bitstream
s = slide(dark=True)
kicker(s, 'Part 1 — What an FPGA is', dark=True)
title(s, 'So what is a "bitstream"?', dark=True)
logo(s, ASSETS + '/lattice_logo_ondark.png', h=0.34)
bullets(s, 0.6, 1.95, 6.6, 4.6, [
    '- Every LUT’s 16 bits. Every routing switch. Every I/O standard. Every block-RAM’s '
    'initial contents.',
    '- All of it is just **configuration SRAM** spread across the die.',
    '- The bitstream is a dump of that SRAM — a few megabits — shifted in over JTAG at '
    'power-up.',
    '- SRAM-based ⇒ **volatile**. Power cycle and the chip forgets it was a CPU. Either '
    'reload it, or store it in the board’s SPI flash.',
], dark=True, size=17, space=16)
chain(s, 7.5, 2.3, ['.bit file', 'JTAG', 'config\nSRAM', 'the chip\nIS the\ncircuit'],
      box_w=1.1, box_h=1.05, gap=0.26, dark=True, size=11)
code_panel(s, 7.5, 3.9, 5.3, 1.5,
           '$ openFPGALoader -c ft2232 -m b8008.bit\n'
           'Jtag frequency : 6.00MHz\n'
           'Loading: [==================] 100.00%\n'
           'Done', size=12, bg=RGBColor(0x26, 0x29, 0x2D), comment=None)
sh = rect(s, 7.5, 5.65, 5.3, 0.95, CHARCOAL2)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'This design’s bitstream: **1.7 Mbit** — the same file, every power-up.',
     True, 16)
notes(s, 'Kill the most common misconception in the room: the bitstream is not a program '
         'and the FPGA is not executing it. The bitstream is the *wiring diagram*, written '
         'into memory cells that hold switches open or closed. After that it is inert — '
         'the circuit just is.')

# ---- Slide 8: Not emulation
s = slide()
kicker(s, 'Part 1 — What an FPGA is')
title(s, 'This is not emulation. Everything happens at once.')
# left: software
rect(s, 0.6, 1.95, 5.85, 4.5, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1.25)
rect(s, 0.6, 1.95, 5.85, 0.09, MIDGREY)
tb, tf = box(s, 0.85, 2.15, 5.3, 0.5)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'An 8008 emulator (software)'
r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True
r.font.color.rgb = CHARCOAL
code_panel(s, 0.85, 2.8, 5.3, 2.05,
           'while True:\n'
           '    op = mem[pc]      # one\n'
           '    pc += 1           # thing\n'
           '    decode(op)        # at a\n'
           '    execute(op)       # time', size=13, comment='#')
bullets(s, 0.85, 5.0, 5.3, 1.3, [
    'Sequential. Fast, but a *model* of the machine — cycle counts are bookkeeping, '
    'not physics.',
], size=15)
# right: fpga
rect(s, 6.85, 1.95, 5.85, 4.5, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=YELLOW, line_w=2)
rect(s, 6.85, 1.95, 5.85, 0.09, YELLOW)
tb, tf = box(s, 7.1, 2.15, 5.3, 0.5)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'An 8008 in an FPGA'
r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True
r.font.color.rgb = CHARCOAL
mods = ['ALU', 'register file', 'stack', 'decoder', 'timing gen', 'flags',
        'temp regs', 'I/O buffer', 'addr mux']
mx, my = 7.1, 2.8
for i, m in enumerate(mods):
    rr, cc = divmod(i, 3)
    sh = rect(s, mx + cc * 1.78, my + rr * 0.62, 1.66, 0.5, PALE_YELLOW,
              line_color=YELLOW, line_w=1)
    shape_text(sh, m, size=11, color=CHARCOAL)
bullets(s, 7.1, 4.85, 5.3, 1.4, [
    'All nine blocks settle **simultaneously**, every clock edge. There is no loop. '
    'There is no "next".',
], size=15)
sh = rect(s, 0.6, 6.6, 12.1, 0.75, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'The board is not *running* an 8008. The board **is** an 8008.', True, 18)
notes(s, 'The parallelism point is the reason people use FPGAs at all — SDR, video, HFT, '
         'protocol bridges. Here it also buys authenticity: the timing is real, so period '
         'software that depends on cycle counts just works. If you want one sentence for '
         'the software folks: in HDL, every line you write runs on every clock edge, '
         'forever, at the same time as every other line.')

# ---- Slide 9: The board
s = slide(dark=True)
kicker(s, 'Part 1 — What an FPGA is', dark=True)
title(s, 'The board I used: Lattice ECP5-5G Versa', dark=True)
logo(s, ASSETS + '/lattice_logo_ondark.png', h=0.34)
bullets(s, 0.6, 1.95, 6.4, 4.3, [
    '- **LFE5UM5G-45F** — 43,848 LUT4/FF, 108 block RAMs (~1.9 Mb), 4 PLLs, 5 Gbps SerDes',
    '- ~$150 dev kit: DIP switches and LEDs become the 8008’s **front panel**; '
    'USB-UART is the terminal',
    '- Chosen for one reason above all: **the open-source tools know this chip cold** '
    '(Project Trellis reverse-engineered the bitstream format)',
    '- No vendor license, no web-install, no node-locked seat. `git clone` and build.',
], dark=True, size=16, space=15)
stat_tiles(s, 7.4, 2.15, [('44K', 'logic cells'), ('108', 'block RAMs')],
           tile_w=2.5, tile_h=1.4, dark=True, num_size=32)
stat_tiles(s, 7.4, 3.75, [('25 MHz', 'system clock'), ('500 kHz', '8008 needs')],
           tile_w=2.5, tile_h=1.4, dark=True, num_size=26)
notes(s, 'The honest reason for ECP5: open toolchain. Xilinx/Intel parts are bigger and '
         'faster but you are inside a closed binary. On ECP5 I could read the synthesis '
         'log, read the packer, and — as you will see — file a bug against Yosys when it '
         'ate my ROM. Optional: hold up the board here.')

# ============================================================
# PART 2 — HDL
# ============================================================

# ---- Slide 10: HDL is not a programming language
s = slide()
kicker(s, 'Part 2 — Describing hardware')
title(s, 'HDL is not a programming language')
bullets(s, 0.6, 1.9, 12.2, 0.6, [
    'It looks like one. It is not. Three differences do all the damage:',
], size=17)
items = [
    ('Everything is concurrent',
     'Statements are not steps. They are **wires that exist all the time**. '
     'Reordering two lines changes nothing.'),
    ('You describe structure,\nnot behavior over time',
     'Your `if` becomes a multiplexer. Your `+` becomes an adder. '
     '**The code is a schematic in text form.**'),
    ('The clock is the only\n"when" you get',
     'Nothing happens "next". Things happen **at a clock edge**, or they '
     'happen continuously.'),
]
cx = 0.6
for head, body in items:
    sh = rect(s, cx, 2.7, 3.9, 3.4, RGBColor(0xFF, 0xFF, 0xFF),
              line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
    rect(s, cx, 2.7, 0.08, 3.4, YELLOW)
    tf2 = sh.text_frame; tf2.word_wrap = True
    tf2.margin_left = Inches(0.28); tf2.margin_right = Inches(0.2)
    tf2.margin_top = Inches(0.25)
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = head
    r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    p = tf2.add_paragraph(); p.space_before = Pt(14); p.alignment = PP_ALIGN.LEFT
    rich(p, body, False, 14, base_color=RGBColor(0x5A, 0x60, 0x68))
    cx += 4.1
sh = rect(s, 0.6, 6.35, 12.1, 0.85, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'If you find yourself asking "what happens first?", you are still writing software.',
     True, 17)
notes(s, 'The single hardest unlearning for a software engineer. My own worst early bugs '
         'were all this: assuming sequence where there was none. A good check when '
         'reviewing HDL — ask "what gate is this line?" If you cannot answer, it probably '
         'will not synthesize the way you think.')

# ---- Slide 11: VHDL vs Verilog
s = slide()
kicker(s, 'Part 2 — Describing hardware')
title(s, 'Two languages, one job')
code_panel(s, 0.6, 2.0, 5.9, 3.5,
           '-- VHDL  (Ada-flavoured)\n'
           'process(clk)\n'
           'begin\n'
           '  if rising_edge(clk) then\n'
           '    if rst = \'1\' then\n'
           '      count <= (others => \'0\');\n'
           '    else\n'
           '      count <= count + 1;\n'
           '    end if;\n'
           '  end if;\n'
           'end process;', size=14)
code_panel(s, 6.85, 2.0, 5.9, 3.5,
           '// Verilog  (C-flavoured)\n'
           'always @(posedge clk)\n'
           'begin\n'
           '  if (rst)\n'
           '    count <= 8\'b0;\n'
           '  else\n'
           '    count <= count + 1;\n'
           'end', size=14, comment='//')
caption(s, 0.6, 5.6, 5.9, 'Verbose. Strongly typed. Refuses to guess.')
caption(s, 6.85, 5.6, 5.9, 'Terse. Weakly typed. Dominant in US industry.')
bullets(s, 0.6, 6.1, 12.2, 1.2, [
    'Same synthesized result: **8 flip-flops and an incrementer on a carry chain.**',
    'I chose VHDL because it will not silently truncate a vector — and when a large share '
    'of your code is drafted by an AI, a compiler that refuses to guess is a feature.',
], size=16, space=8)
notes(s, 'Do not turn this into a language war — 60 seconds. The point is that HDLs are '
         'interchangeable at this level, and that the *type strictness* of VHDL was a '
         'deliberate safety choice for this particular project. Modern tools mix both '
         'freely: this design is VHDL that gets translated to Verilog on the way to the '
         'chip anyway.')

# ---- Slide 12: The three things you write
s = slide()
kicker(s, 'Part 2 — Describing hardware')
title(s, 'You only ever write three kinds of thing')
col_w = 3.8
heads = [
    ('1 — Combinational', 'wires and gates. No memory.', GREEN),
    ('2 — Sequential', 'flip-flops. Remembers.', SKY),
    ('3 — Structural', 'wiring sub-blocks together.', YELLOW),
]
codes = [
    'sum <= a xor b xor cin;\n'
    'cout <= (a and b)\n'
    '     or (cin and (a xor b));\n'
    '\n'
    '-- exists continuously.\n'
    '-- becomes: 2 LUT4s.',

    'process(clk)\n'
    'begin\n'
    '  if rising_edge(clk) then\n'
    '    if push = \'1\' then\n'
    '      sp <= sp + 1;\n'
    '    end if;\n'
    '  end if;\n'
    'end process;\n'
    '-- becomes: 3 TRELLIS_FF.',

    'u_alu : entity work.alu\n'
    '  port map (\n'
    '    op_a  => temp_a,\n'
    '    op_b  => temp_b,\n'
    '    result=> alu_out,\n'
    '    carry => flag_c\n'
    '  );\n'
    '\n'
    '-- becomes: the netlist.',
]
cx = 0.6
for (head, sub, color), code in zip(heads, codes):
    rect(s, cx, 1.95, col_w, 0.09, color)
    tb, tf = box(s, cx, 2.1, col_w, 0.85)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = head
    r.font.name = FONT; r.font.size = Pt(19); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    p = tf.add_paragraph()
    r = p.add_run(); r.text = sub
    r.font.name = FONT; r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0x5A, 0x60, 0x68)
    code_panel(s, cx, 3.05, col_w, 2.9, code, size=12.5)
    cx += col_w + 0.35   # 0.6 / 4.75 / 8.90 -> right edge 12.70
bullets(s, 0.6, 6.15, 12.2, 1.1, [
    'That is the whole language, for synthesis purposes. Everything else — types, '
    'generics, packages — is bookkeeping so humans can manage thousands of these.',
    'This design: **29 files**, each one mostly kind 1 and 2, wired by one file of kind 3.',
], size=16, space=8)
notes(s, 'Demystifying move: HDL has a small synthesizable core. Most of a VHDL textbook '
         'is simulation-only constructs that never reach the chip. If you can read these '
         'three patterns you can read this entire project. Note the comments: I want the '
         'audience already thinking in "what does it become".')

# ---- Slide 13: A whole real module
s = slide()
kicker(s, 'Part 2 — Describing hardware')
title(s, 'A whole real module, start to finish')
code_panel(s, 0.6, 1.95, 7.0, 4.5,
           'entity stack_pointer is\n'
           '  port ( clk, reset   : in  std_logic;\n'
           '         phi1_rising  : in  std_logic;\n'
           '         stack_push   : in  std_logic;\n'
           '         stack_pop    : in  std_logic;\n'
           '         sp           : out unsigned(2 downto 0) );\n'
           'end entity;\n'
           '\n'
           'process(clk, reset) begin\n'
           '  if reset = \'1\' then\n'
           '    sp_r <= (others => \'0\');\n'
           '  elsif rising_edge(clk) then\n'
           '    if phi1_rising = \'1\' then\n'
           '      if    stack_push = \'1\' then sp_r <= sp_r + 1;\n'
           '      elsif stack_pop  = \'1\' then sp_r <= sp_r - 1;\n'
           '      end if;\n'
           '    end if;\n'
           '  end if;\n'
           'end process;', size=13)
bullets(s, 7.95, 1.95, 4.85, 4.6, [
    'This is the **entire** 8008 stack pointer. 69 lines with comments.',
    'It does not know what `CALL` is. It does not know interrupts exist. It counts, and it '
    'wraps 7→0 for free because it is 3 bits wide.',
    'In silicon: **3 flip-flops, an incrementer, a decrementer, and a mux.**',
    'It has never had a bug.',
], size=16, space=16)
sh = rect(s, 0.6, 6.55, 12.1, 0.75, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'Design rule for the whole project: **every module is dumb.** '
        'One job, ~50–100 lines, no knowledge of any other module.', True, 17)
notes(s, 'Two payoffs, both worth naming. (1) Verification: a module this small is '
         'exhaustively testable. (2) It is the unit an AI assistant gets right and a human '
         'can fully review — which mattered, because most of this VHDL was drafted by '
         'Claude and reviewed by me. The 8008’s wrap-around stack semantics fall out of '
         'the 3-bit width; I never wrote code for it.')

# ---- Slide 14: VHDL to LUTs, all the way down
s = slide()
kicker(s, 'Part 2 — Describing hardware')
title(s, 'All the way down: one line of VHDL → silicon')
stages = [
    ('You write', 'sp_r <= sp_r + 1;', '-- VHDL'),
    ('GHDL emits', 'assign n42 = sp_r + 3\'d1;', '// Verilog'),
    ('Yosys maps', 'CCU2C  #(.INIT0(...))\n  ccu_0 (.A0(sp_r[0]), ...);', '// ECP5 cells'),
    ('nextpnr places', 'X23/Y7/SLICEA  →  X23/Y7/SLICEB', '// coordinates'),
]
cy = 1.95
for i, (who, code, tag) in enumerate(stages):
    sh = rect(s, 0.6, cy, 2.5, 0.85, CHARCOAL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(sh, who, size=15, bold=True)
    code_panel(s, 3.35, cy, 6.2, 0.85, code, size=12.5,
               comment='--' if i == 0 else '//')
    tb, tf = box(s, 9.75, cy + 0.14, 3.1, 0.7)
    p = tf.paragraphs[0]
    rich(p, tag.replace('--', '').replace('//', '').strip(), False, 13,
         base_color=MIDGREY)
    if i < len(stages) - 1:
        arrow(s, 1.85, cy + 0.87, 1.85, cy + 1.03, color=YELLOW, wpt=2.5)
    cy += 1.05
sh = rect(s, 0.6, cy + 0.05, 12.1, 0.95, PALE_YELLOW, line_color=YELLOW, line_w=1.5)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'Every intermediate file is on my disk, in plain text, because the whole '
        'toolchain is open source.', False, 16)
notes(s, 'This is the spine of the talk for hardware-curious folks: the abstraction stack '
         'is only four levels deep and you can inspect all of it. Mention that '
         '`build/*.v` and `build/*.json` are readable — I have debugged real problems by '
         'grepping the Yosys JSON. Contrast with a closed vendor flow where the '
         'intermediate is a binary blob.')

# ---- Slide 15: Inference
s = slide()
kicker(s, 'Part 2 — Describing hardware')
title(s, 'Inference: say the magic words, get the hard block')
code_panel(s, 0.6, 1.95, 6.5, 2.65,
           'process(CLK)\n'
           'begin\n'
           '  if rising_edge(CLK) then\n'
           '    if CS_N=\'0\' and RW_N=\'0\' then\n'
           '      ram(idx) <= DATA_IN;\n'
           '    end if;\n'
           '    DATA_OUT <= ram(idx);   -- registered read\n'
           '  end if;\n'
           'end process;', size=13)
caption(s, 0.6, 4.7, 6.5, 'src/b8008/ram_sync.vhdl — the 8008’s memory')
bullets(s, 7.4, 1.95, 5.4, 2.8, [
    'Because **both** the read and the write are clocked, Yosys recognizes the shape and '
    'emits `DP16KD` — a real block RAM.',
    'All 16 KB of this machine’s ROM and RAM: **8 block RAMs. Zero LUTs.**',
], size=16, space=14)
# right/wrong panel
rect(s, 0.6, 5.15, 6.5, 1.95, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=RED, line_w=1.5)
tb, tf = box(s, 0.85, 5.28, 6.0, 1.7)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'Move the read outside the clocked process…'
r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
r.font.color.rgb = RED
p = tf.add_paragraph(); p.space_before = Pt(8)
rich(p, '…and you have asked for an **asynchronous-read memory**, which a block RAM '
        'physically cannot do. Yosys obliges — by building it out of **131,072 '
        'flip-flops**. The chip has 43,848. It does not fit, and the error arrives three '
        'tools downstream, phrased as a placement failure.', False, 14)
rect(s, 7.4, 5.15, 5.4, 1.95, PALE_YELLOW, line_color=YELLOW, line_w=1.5)
tb, tf = box(s, 7.62, 5.3, 5.0, 1.7)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'The lesson'
r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
r.font.color.rgb = CHARCOAL
p = tf.add_paragraph(); p.space_before = Pt(8)
rich(p, 'In HDL, **the shape of your code is the request.** Same logic, same result in '
        'simulation — one version fits in 7% of the block RAM, the other does not fit '
        'at all.', False, 14)
notes(s, 'Best single example of "HDL is not software". Both versions simulate '
         'identically and both are correct. One fits in 7% of the block RAM; the other '
         'does not fit at all. This is why HDL people talk about "coding styles" — they '
         'are not style, they are load-bearing. Also note the deliberate design comment in '
         'ram_sync.vhdl: the one-cycle read latency is invisible to the 8008 because the '
         'address is latched at T1/T2 and data is not consumed until T3.')

# ---- Slide 16: Things that don't exist inside an FPGA
s = slide()
kicker(s, 'Part 2 — Describing hardware')
title(s, 'Things that do not exist inside an FPGA')
rows = [
    ('Tri-state buses',
     'The 8008 has one 8-bit bus doing three jobs. Inside the fabric there are **no '
     'internal tri-states** — only at the I/O pins.',
     'Yosys `tribuf -logic` rewrites every `\'Z\'` into a multiplexer. My "bus" is a mux '
     'tree, and it behaves identically.'),
    ('Free initial values',
     'A signal is not `0` at power-up because you said so. Block RAM contents come from '
     'the bitstream; logic comes from your reset.',
     'If you rely on an initial value that the fabric will not honour, simulation is '
     'right and hardware is wrong. Reset everything you care about.'),
    ('"Wait a moment"',
     'There is no delay primitive. `after 10 ns` is a **simulation-only** construct and '
     'synthesizes to nothing.',
     'Time is counted in clock cycles. My 0.8 µs φ1 pulse is a counter comparing against '
     '20 ticks of a 25 MHz clock.'),
]
cy = 1.95
for head, what, fix in rows:
    rect(s, 0.6, cy, 12.1, 1.5, RGBColor(0xFF, 0xFF, 0xFF),
         line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
    rect(s, 0.6, cy, 0.09, 1.5, YELLOW)
    tb, tf = box(s, 0.85, cy + 0.14, 2.9, 1.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = head
    r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    tb, tf = box(s, 3.9, cy + 0.12, 4.3, 1.3)
    p = tf.paragraphs[0]
    rich(p, what, False, 13.5, base_color=RGBColor(0x4A, 0x50, 0x58))
    tb, tf = box(s, 8.4, cy + 0.12, 4.1, 1.3)
    p = tf.paragraphs[0]
    rich(p, fix, False, 13.5, base_color=RGBColor(0x4A, 0x50, 0x58))
    cy += 1.66
tb, tf = box(s, 3.9, 1.62, 4.3, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'WHAT BITES'
r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = MIDGREY
tb, tf = box(s, 8.4, 1.62, 4.1, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'WHAT YOU DO'
r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = MIDGREY
notes(s, 'These three are where a 1972 datasheet collides with 2013 silicon. The tri-state '
         'one is the fun one: the 8008’s entire personality is a time-multiplexed bus, and '
         'the FPGA simply does not have that primitive internally — so the tool rebuilds '
         'the same behavior out of muxes. Faithful function, completely different '
         'implementation.')

# ============================================================
# PART 3 — THE TOOLCHAIN
# ============================================================

# ---- Slide 17: The flow
s = slide(dark=True)
kicker(s, 'Part 3 — The toolchain', dark=True)
title(s, 'Source to silicon, entirely open source', dark=True)
logo(s, ASSETS + '/lattice_logo_ondark.png', h=0.34)
FLOW_W, FLOW_GAP = 1.62, 0.40   # 0.6 + 6*W + 5*GAP = 12.32
chain(s, 0.6, 2.15, ['VHDL', 'GHDL', 'Yosys', 'nextpnr', 'ecppack', 'board'],
      box_w=FLOW_W, box_h=0.95, gap=FLOW_GAP, dark=True, size=14)
labels = ['29 files', 'analyze +\nemit Verilog', 'synthesize to\nECP5 cells',
          'place & route,\ntiming', 'pack .bit', 'openFPGA\nLoader']
cx = 0.6
for lbl in labels:
    tb, tf = box(s, cx - 0.14, 3.18, FLOW_W + 0.28, 0.8)
    for i, ln in enumerate(lbl.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(11.5)
        r.font.color.rgb = RGBColor(0xB6, 0xBC, 0xC4)
    cx += FLOW_W + FLOW_GAP
code_panel(s, 0.6, 4.15, 12.1, 2.15,
           '$ ghdl --synth --std=08 --out=verilog b8008_basic_top  > build/top.v\n'
           '$ yosys -p "read_verilog src/synth/ghdl_gates.v build/top.v;  \\\n'
           '            tribuf -logic; proc; opt -nodffe;                 \\\n'
           '            synth_ecp5 -top b8008_basic_top -json build/top.json"\n'
           '$ nextpnr-ecp5 --um5g-45k --package CABGA381 --speed 8 \\\n'
           '            --json build/top.json --lpf constraints/b8008_basic.lpf \\\n'
           '            --textcfg build/top.config\n'
           '$ ecppack --input build/top.config --bit build/top.bit',
           size=12.5, bg=RGBColor(0x26, 0x29, 0x2D), comment=None)
bullets(s, 0.6, 6.5, 12.1, 0.8, [
    'Four tools, one Makefile, **no vendor license**. Full rebuild: about 90 seconds.',
], dark=True, size=16)
notes(s, 'These are the real commands out of projects/project.mk. Worth calling out '
         '`tribuf -logic` — that is the tri-state rewrite from two slides ago, sitting '
         'right there in the flow. And `ghdl_gates.v` is a 40-line shim I had to write '
         'because GHDL emits two primitives Yosys does not know. Both are the kind of '
         'seam you only get to see, let alone fix, in an open toolchain.')

# ---- Slide 18: What P&R actually does
s = slide()
kicker(s, 'Part 3 — The toolchain')
title(s, 'The two hard steps: place, then route')
rect(s, 0.6, 1.95, 5.9, 3.6, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1.25)
rect(s, 0.6, 1.95, 5.9, 0.09, YELLOW)
tb, tf = box(s, 0.85, 2.12, 5.4, 0.5)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'PLACE'
r.font.name = FONT; r.font.size = Pt(20); r.font.bold = True
r.font.color.rgb = CHARCOAL
bullets(s, 0.85, 2.7, 5.4, 2.7, [
    'Assign each of ~2,400 cells to a physical site on the die.',
    'Objective: things that talk to each other should sit near each other.',
    'It is simulated annealing. It is **not deterministic across versions**, and two '
    'builds can differ in fmax by 20%.',
], size=15, space=12)
rect(s, 6.85, 1.95, 5.9, 3.6, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1.25)
rect(s, 6.85, 1.95, 5.9, 0.09, YELLOW)
tb, tf = box(s, 7.1, 2.12, 5.4, 0.5)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'ROUTE'
r.font.name = FONT; r.font.size = Pt(20); r.font.bold = True
r.font.color.rgb = CHARCOAL
bullets(s, 7.1, 2.7, 5.4, 2.7, [
    'Find a path through the switch matrix for every net.',
    'Wires are a **finite resource**. A design can fit logically and still fail to route.',
    'Every hop costs picoseconds — which is why placement quality decides your clock '
    'speed, not gate count.',
], size=15, space=12)
sh = rect(s, 0.6, 5.75, 12.1, 1.45, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'Software analogy: synthesis is the compiler, place & route is the linker — '
        'except this linker also decides how fast your program runs, and takes minutes '
        'to do it.', True, 18)
notes(s, 'Manage expectations for anyone who might try this: the edit-compile-test loop '
         'is minutes, not seconds, and it is stochastic. That changes how you work — you '
         'batch changes, you lean much harder on simulation, and you do not chase a 3% '
         'fmax difference between two runs because it is noise.')

# ---- Slide 19: Timing closure
s = slide()
kicker(s, 'Part 3 — The toolchain')
title(s, 'Timing closure: the only question that matters')
bullets(s, 0.6, 1.9, 12.2, 0.8, [
    'Between two clock edges, every signal must finish traveling. **Longest path wins.** '
    'That path is your maximum clock frequency.',
], size=17)
code_panel(s, 0.6, 2.85, 7.6, 2.7,
           '        type  curr total\n'
           '    clk-to-q  1.36  1.36  ROM block RAM .DOA0\n'
           '     routing  1.18  2.54  (42,10) -> (31,6)\n'
           '       logic  0.24  2.78  cpu_s0 decode\n'
           '     routing  0.86  3.64  (31,6)  -> (28,3)\n'
           '           .. 14 more hops ..\n'
           '     routing  1.13  8.83  cpu_d[0] -> pin\n'
           '\n'
           'Info: 3.09 ns logic, 5.74 ns routing\n'
           'Info: Max frequency: 115.66 MHz (PASS at 25.00 MHz)',
           size=13, comment=None)
caption(s, 0.6, 5.65, 7.6,
        'Real output, abridged — projects/b8008_basic/reports/pnr.txt')
bullets(s, 8.5, 2.85, 4.3, 3.4, [
    'nextpnr does the split for you: **3.09 ns of logic, 5.74 ns of wire.**',
    'Routing is nearly **two-thirds** of the delay. That is the norm, not an anomaly.',
    'We need 25 MHz and close at 115.66 — **4.6× margin**, so I never had to fight for it.',
    'If it had failed: pipeline the path, re-floorplan, or slow the clock.',
], size=15, space=13)
sh = rect(s, 0.6, 6.35, 12.1, 0.85, PALE_YELLOW, line_color=YELLOW, line_w=1.5)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'The 8008 itself runs at **500 kHz**. I had 230× the speed I needed — and I still '
        'ran the fabric at 25 MHz just to divide down cleanly.', False, 16)
notes(s, 'Two things to land. (1) What "fmax" means physically — it is a race between a '
         'signal and the next clock edge. (2) This path starts in the ROM block RAM, '
         'wanders through instruction decode and the register file, and ends at a data-bus '
         'pin — 24 hops, and the routing between tiles costs more than all the logic. '
         'Note for honesty: nextpnr prints fmax twice, once after placement (94.58 MHz) '
         'and again after routing. 115.66 is the final one. If someone asks how you fix a '
         'failing path: add a pipeline register and let the tool re-place it.')

# ---- Slide 20: Utilization
s = slide()
kicker(s, 'Part 3 — The toolchain')
title(s, 'What a 1972 computer costs in 2013 silicon')
table(s, 0.6, 1.95, 7.4, 3.3, [
    ['Resource', 'Used', 'Available', '%'],
    ['TRELLIS_COMB (LUTs)', '1,770', '43,848', '**4%**'],
    ['TRELLIS_FF (flip-flops)', '591', '43,848', '1%'],
    ['DP16KD (block RAM)', '8', '108', '7%'],
    ['EHXPLLL (PLL)', '1', '4', '25%'],
    ['TRELLIS_IO (pins)', '63', '245', '25%'],
], [2.9, 1.5, 1.6, 1.4], size=15)
caption(s, 0.6, 5.35, 7.4,
        'b8008_basic: CPU + 12 KB ROM + 4 KB RAM + UART + clock generation + front panel')
stat_tiles(s, 8.4, 1.95, [('4%', 'of the fabric')], tile_w=4.3, tile_h=1.55,
           num_size=54, lbl_size=14)
bullets(s, 8.4, 3.75, 4.4, 3.0, [
    'The **whole system** — CPU, memory, serial port, clock generator — fits in one '
    'twenty-fifth of a mid-range 2013 FPGA.',
    'The 8008 was 3,500 transistors. This chip has roughly **a million times** the logic '
    'and it is not an expensive part.',
    'There is room on this die for 25 of these machines, running side by side.',
], size=15, space=14)
notes(s, 'Good "wow" beat, and an honest one — those numbers are straight from '
         'reports/utilization.txt. The follow-up question you will get is "so what would '
         'you do with the other 96%?" Real answer: a 6502, a Z80, and a video generator, '
         'all running at once, sharing nothing. That is what the parallelism slide meant '
         'in practice.')

# ============================================================
# PART 4 — THE DESIGN + WAR STORIES
# ============================================================

# ---- Slide 22: Where the architecture came from (NEW)
s = slide()
kicker(s, 'Part 4 — Building the thing')
title(s, 'Where the architecture came from')
bw, bh = fit(ASSETS + '/8008_block_diagram.png', 5.5, 2.65)
pic(s, ASSETS + '/8008_block_diagram.png', 0.6, 1.9, w=bw)
caption(s, 0.6, 1.9 + bh + 0.06, 5.5,
        'Intel 8008 User’s Manual, 1972 — the block diagram')
bullets(s, 6.4, 1.9, 6.3, 2.9, [
    'Every module in this design is **one box from that diagram.**',
    'Intel drew the *concept*. The RTL partition is mine: **29 independently testable '
    'modules**, each with explicit control signals and no knowledge of any other — '
    'derived from the datasheet, not from any existing implementation.',
], size=16, space=16)
# the claim
sh = rect(s, 0.6, 4.95, 12.1, 1.85, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.35); tf.margin_right = Inches(0.35)
tf.margin_top = Inches(0.18)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'As far as I can tell, this is the only block-based 8008 in existence.'
r.font.name = FONT; r.font.size = Pt(21); r.font.bold = True
r.font.color.rgb = YELLOW
p = tf.add_paragraph(); p.space_before = Pt(9); p.alignment = PP_ALIGN.CENTER
rich(p, 'Eighteen months of looking — Google Groups, websites that look like they are '
        'from 1994, the SCELBI archives. Everything I found, and everything in the '
        'model’s training data, is an **instruction-set simulator**: sequential software '
        'that executes one opcode at a time.', True, 14)
p = tf.add_paragraph(); p.space_before = Pt(7); p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'If you find another, please tell me. I would genuinely like to see it.'
r.font.name = FONT; r.font.size = Pt(14); r.font.italic = True
r.font.color.rgb = RGBColor(0xC9, 0xCE, 0xD6)
notes(s, 'This is the keystone slide and it is worth slowing down for. Two separate '
         'claims, kept separate. First, provenance: Intel drew a conceptual block '
         'diagram; turning those boxes into 29 synthesizable, independently testable '
         'modules with explicit control signals is an RTL partition, and that work is '
         'mine. Second, the survey: eighteen months of searching turned up only '
         'instruction-set simulators. Make the invitation genuinely — you want the '
         'counterexample if it exists. The corpus observation is not a throwaway; it is '
         'the mechanism behind the two failed attempts, and you cash it in at the end.')

# ---- Slide 21: What is actually on the chip
s = slide()
kicker(s, 'Part 4 — Building the thing')
title(s, 'The whole computer, on one die')
# CPU core box
rect(s, 0.6, 1.95, 5.4, 4.3, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=CHARCOAL, line_w=2)
tb, tf = box(s, 0.8, 2.05, 5.0, 0.45)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'b8008 CPU core — 29 VHDL modules'
r.font.name = FONT; r.font.size = Pt(15); r.font.bold = True
r.font.color.rgb = CHARCOAL
core = ['timing &\ncontrol', 'instruction\ndecoder', 'address\nstack (8×14)',
        'stack\npointer', 'register file\nA B C D E H L', 'ALU +\ncarry lookahead',
        'condition\nflags', 'temp\nregisters', 'I/O buffer']
mx, my = 0.85, 2.6
for i, m in enumerate(core):
    rr, cc = divmod(i, 3)
    sh = rect(s, mx + cc * 1.72, my + rr * 1.15, 1.6, 1.0, PALE_YELLOW,
              line_color=YELLOW, line_w=1)
    shape_text(sh, m, size=11, color=CHARCOAL)
# peripherals
per = [
    ('rom_12kx8_bram', '12 KB firmware — 6 × DP16KD'),
    ('ram_sync', '4 KB program RAM — 2 × DP16KD'),
    ('address_decoder', 'RAM 0000-0FFF, ROM 1000-3FFF'),
    ('usart', '115200 8N1 over USB'),
    ('phase_clocks', 'φ1 / φ2, 0.8 µs / 0.6 µs, non-overlapping'),
    ('EHXPLLL', 'hard PLL: 100 MHz board → 25 MHz system'),
    ('debug_clock_control', 'freeze / single-step the CPU'),
]
py = 1.95
for name, desc in per:
    sh = rect(s, 6.4, py, 6.3, 0.56, RGBColor(0xFF, 0xFF, 0xFF),
              line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
    tf2 = sh.text_frame; tf2.word_wrap = True
    tf2.margin_left = Inches(0.15); tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = name
    r.font.name = MONO; r.font.size = Pt(12.5); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    r = p.add_run(); r.text = '   ' + desc
    r.font.name = FONT; r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0x6A, 0x70, 0x78)
    py += 0.62
sh = rect(s, 0.6, 6.45, 12.1, 0.8, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'In 1972 all of this was a rack. Intel’s own SIM8-01 reference board needed dozens '
        'of TTL chips just to make the CPU usable.', True, 16)
notes(s, 'This is the "system on chip" slide. The point for an FPGA audience: the CPU is '
         'the small part. Memory, clocking, serial, and debug infrastructure are most of '
         'the engineering — same as any real SoC. Call out debug_clock_control: being able '
         'to freeze the CPU with a DIP switch and single-step it is a debugging superpower '
         'you simply cannot have with a physical 1972 chip.')

# ---- Slide 22: Clocking trap
s = slide()
kicker(s, 'Part 4 — Building the thing')
title(s, 'The trap I fell in: φ1 and φ2 are data, not clocks')
bullets(s, 0.6, 1.9, 6.3, 2.3, [
    'The 8008 needs **two non-overlapping clock phases**. The obvious move is to generate '
    'them and clock your logic from them.',
    'Do that and you have just created **two new clock domains** — hand-made, skewed, '
    'un-analyzable by the timing engine, and impossible to close.',
], size=16, space=14)
rect(s, 0.6, 4.3, 6.3, 2.5, RGBColor(0xFF, 0xFF, 0xFF),
     line_color=YELLOW, line_w=2)
rect(s, 0.6, 4.3, 6.3, 0.09, YELLOW)
tb, tf = box(s, 0.85, 4.45, 5.8, 2.2)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'The fix'
r.font.name = FONT; r.font.size = Pt(17); r.font.bold = True
r.font.color.rgb = CHARCOAL
p = tf.add_paragraph(); p.space_before = Pt(8)
rich(p, 'One real clock — 25 MHz from the PLL. φ1 and φ2 become **ordinary signals**, '
        'plus one-cycle `phi1_rising` / `phi2_rising` enable pulses. Every flip-flop in '
        'the design is clocked by `clk_sys` and gated by an enable.', False, 14.5)
p = tf.add_paragraph(); p.space_before = Pt(8)
rich(p, 'Result: **one clock domain**, fully timed, closes at 115 MHz.', False, 14.5)
code_panel(s, 7.25, 1.95, 5.5, 2.65,
           '-- WRONG: phi1 as a clock\n'
           'process(phi1)\n'
           'begin\n'
           '  if rising_edge(phi1) then\n'
           '    sp <= sp + 1;\n'
           '  end if;\n'
           'end process;', size=13)
code_panel(s, 7.25, 4.85, 5.5, 1.95,
           '-- RIGHT: phi1 as an enable\n'
           'if rising_edge(clk) then\n'
           '  if phi1_rising = \'1\' then\n'
           '    sp <= sp + 1;\n'
           '  end if;\n'
           'end if;', size=13)
notes(s, 'Best transferable FPGA lesson in the deck. The general rule: derive clocks in a '
         'PLL or not at all; everything else is an enable. Divided/gated clocks are the '
         'classic beginner mistake and they produce bugs that only show up on hardware, '
         'because simulation has no clock skew. The scope captures on the next-to-last '
         'slide show the real φ1/φ2 coming out of the FPGA — those are outputs, not '
         'internal clocks.')

# ---- Slide 23: The 0xFF bug
s = slide(dark=True)
kicker(s, 'Part 4 — Building the thing', dark=True)
title(s, 'War story: the day the toolchain ate my ROM', dark=True)
steps = [
    ('Symptom', 'CPU hangs at boot on hardware. **Simulation passes perfectly.**'),
    ('Red herring #1', 'Adding 32 bytes of padding "fixes" it. 399 bytes → crash; '
                       '431 bytes → fine. Must be a size threshold?'),
    ('Red herring #2', 'Shrink the ROM to 1 KB → *different* failure. Runs, but no UART '
                       'output at all. So, not size.'),
    ('Root cause', 'The ROM array was initialized `(others => x"FF")`. Yosys’ optimizer '
                   'sees a mostly-all-ones memory and mangles the initialization.'),
    ('Fix', 'Change the fill to `x"00"`. **One character.** Works perfectly, and has '
            'ever since.'),
]
cy = 1.95
for i, (head, body) in enumerate(steps):
    col = YELLOW if i == len(steps) - 1 else CHARCOAL2
    sh = rect(s, 0.6, cy, 2.4, 0.8, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(sh, head, size=14, bold=True,
               color=CHARCOAL if i == len(steps) - 1 else SOFTWHITE)
    tb, tf = box(s, 3.2, cy + 0.02, 9.5, 0.85)
    p = tf.paragraphs[0]
    rich(p, body, True, 14.5)
    cy += 0.92
sh = rect(s, 0.6, cy + 0.05, 12.1, 0.75, CHARCOAL2)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'Written up as a bug report against Yosys — `docs/bug-report-yosys-rom-0xff.md`. '
        'You can only do that when the tools are open.', True, 15)
notes(s, 'Tell this one as a detective story; it is the most memorable slide in the deck '
         'for a technical audience. The moral is not "Yosys is bad" — it is that below '
         'your RTL there is a compiler with bugs, and when simulation and hardware '
         'disagree, the toolchain is a legitimate suspect. Days lost: about three. '
         'Characters changed: one.')

# ---- Slide 24: ecpbram
s = slide()
kicker(s, 'Part 4 — Building the thing')
title(s, 'A trick worth stealing: patch firmware, skip the build')
bullets(s, 0.6, 1.9, 12.2, 0.9, [
    'Problem: my firmware lives in block RAM. Changing one byte of assembly meant a full '
    '90-second resynthesis — and a new place & route, so new timing, for a change that '
    'touches no logic at all.',
], size=17)
chain(s, 0.6, 3.0, ['edit .asm', 'assemble', 'ecpbram\npatch', 'repack .bit', 'flash'],
      box_w=2.0, box_h=1.05, gap=0.4, size=13)
bullets(s, 0.6, 4.5, 6.1, 2.4, [
    'The block-RAM contents sit in the post-route **config file**, not the netlist.',
    '`ecpbram` finds the old contents by pattern-matching and substitutes the new ones.',
    'Unused ROM space is filled with **random bytes** so the match pattern is unique.',
], size=15, space=13)
rect(s, 7.1, 4.5, 5.6, 2.4, PALE_YELLOW, line_color=YELLOW, line_w=1.5)
tb, tf = box(s, 7.35, 4.7, 5.1, 2.1)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = '90 s  →  3 s'
r.font.name = FONT; r.font.size = Pt(38); r.font.bold = True
r.font.color.rgb = CHARCOAL
p = tf.add_paragraph(); p.space_before = Pt(10); p.alignment = PP_ALIGN.CENTER
rich(p, 'Same bitstream, same timing, new firmware. The RTL never gets touched, so it '
        'cannot regress.', False, 14, base_color=RGBColor(0x4A, 0x50, 0x58))
notes(s, 'Generalizable well beyond this project: if your design has a soft CPU or a '
         'lookup table in BRAM, you almost certainly do not need to resynthesize to change '
         'its contents. Every vendor has an equivalent (Xilinx updatemem, Intel '
         'mif-update). It turns a 90-second loop into an interactive one, which changes '
         'how you write firmware.')

# ---- Slide 27: Verification
s = slide()
kicker(s, 'Part 5 — Proof')
title(s, 'How you know gateware is right')
bullets(s, 0.6, 1.9, 12.2, 0.75, [
    'You cannot printf a circuit. Verification is not a phase here — it is the medium '
    'you work in.',
], size=17)
table(s, 0.6, 2.75, 7.5, 3.3, [
    ['Layer', 'What it proves', 'Result'],
    ['Per-module testbench', 'Each dumb block in isolation', '29 modules'],
    ['Exhaustive ALU sweep', 'All 656,384 arithmetic cases vs a reference model', '**PASS**'],
    ['Regression suite', 'Real assembly programs in simulation', '**28 / 28**'],
    ['Cycle-exact T-states', 'Every timing class vs the 1972 datasheet', '**27 / 27**'],
    ['ISA self-test **on the board**', 'The final word', '**46 / 46**'],
], [2.9, 3.4, 1.2], size=14)
rect(s, 8.4, 2.75, 4.3, 3.3, CHARCOAL)
tb, tf = box(s, 8.65, 2.95, 3.8, 3.0)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'The rule I now\nlive by'
r.font.name = FONT; r.font.size = Pt(20); r.font.bold = True
r.font.color.rgb = YELLOW
p = tf.add_paragraph(); p.space_before = Pt(16)
rich(p, 'Simulation passing is a **hypothesis**.', True, 16)
p = tf.add_paragraph(); p.space_before = Pt(10)
rich(p, 'Hardware is the **proof**.', True, 16)
p = tf.add_paragraph(); p.space_before = Pt(16)
rich(p, 'My first attempt passed its entire simulation suite and produced garbage from the '
        'ALU on real silicon. Nothing is done until the board says so.', True, 13.5)
notes(s, 'Bridge from the war stories to the demo. The exhaustive ALU sweep is worth a '
         'beat — 656,384 cases is every operand pair × every op × carry-in, checked '
         'against an independent model. You can do that for an 8-bit ALU; that is a real '
         'advantage of working at this scale. Then: sim-versus-silicon, which is the '
         'lesson every single one of these war stories converges on.')

# ---- Slide 26: Live demo
s = slide(dark=True)
kicker(s, 'Live', dark=True)
title(s, 'The circuit, running', dark=True)
logo(s, ASSETS + '/lattice_logo_ondark.png', h=0.34)
bullets(s, 0.6, 2.0, 7.0, 4.4, [
    '- **The clocks** — φ1 / φ2 on the scope, generated by the fabric',
    '- **Pi to 50 digits** — multiprecision arithmetic, ~75 s at 500 kHz',
    '- **SCELBI floating-point calculator (1974)** — `12.2 X 2.2 = 26.84`',
    '- **The finale** — power on → BASIC → `MON` → machine monitor → `G 1FB6` → back in '
    'BASIC, program intact',
], dark=True, size=17, space=18)
if have(REPO + '/docs/two_phase_oscope_cap_1.png'):
    ow, oh = fit(REPO + '/docs/two_phase_oscope_cap_1.png', 4.9, 2.0)
    pic(s, REPO + '/docs/two_phase_oscope_cap_1.png', 7.9, 2.0, w=ow)
    caption(s, 7.9, 2.0 + oh + 0.08, 4.9,
            'φ1 / φ2 from the FPGA — non-overlapping, 0.8 µs and 0.6 µs',
            color=RGBColor(0x9A, 0xA0, 0xA8))
if have(REPO + '/docs/images/tiny_os_scelbal.png'):
    tw, th = fit(REPO + '/docs/images/tiny_os_scelbal.png', 4.9, 2.4)
    pic(s, REPO + '/docs/images/tiny_os_scelbal.png', 7.9, 4.55, w=tw)
notes(s, 'Demo cheat sheet — verify before the talk. Scope: probe the phi1_out/phi2_out '
         'pins so the audience sees Part 1 made physical. Pi: L in monitor, send_hex.py, '
         'G 0040; narrate while it grinds; 49/50 digits — the 50th is the original 1970s '
         'program’s guard-byte truncation, good trivia. Calc: load, G, type the sum live. '
         'Finale: power-cycle into the b8008_basic bitstream, type a 3-line FOR/NEXT, RUN, '
         'MON, D the program region, G 1FB6, LIST. Fallback recordings ready for all four.')

# ============================================================
# PART 6 — HOW THIS GOT BUILT  (the retrospective sweep)
# ============================================================

# ---- Slide 29: The path — three versions
s = slide(dark=True)
kicker(s, 'Part 6 — How this got built', dark=True)
title(s, 'Three versions, eighteen months', dark=True)
bullets(s, 0.6, 1.5, 12.1, 0.45, [
    'Eighteen months ago I could not read a VHDL `process`, and had never touched an FPGA.',
], dark=True, size=16)
stages = [
    ('STEP ZERO', 'The book',
     'Getting Started with FPGAs —\nRussell Merrick, No Starch Press.',
     'Cover to cover before I wrote\na line of VHDL.', None),
    ('VERSION 1', 's8008',
     'Monolithic, single-cycle.\n1,500 lines in one file.',
     'Passed its entire simulation\nsuite. ALU returned garbage\non the board.', False),
    ('VERSION 2', 'v8008',
     'Multi-cycle rewrite.\nToo clever, too entangled.',
     'Collapsed under its own\ncomplexity.', False),
    ('VERSION 3', 'b8008',
     '29 dumb modules, derived from\nthe 1972 block diagram.',
     'Silicon validated.\n46/46 on the board.', True),
]
cx = 0.6
for label, name, what, outcome, ok in stages:
    sh = rect(s, cx, 2.08, 2.8, 3.25, CHARCOAL2)
    rect(s, cx, 2.08, 2.8, 0.09, YELLOW if ok or ok is None else RGBColor(0x9B, 0x4A, 0x45))
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.24)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(10.5); r.font.bold = True
    r.font.color.rgb = YELLOW if ok or ok is None else RGBColor(0xE0, 0x9A, 0x95)
    p = tf.add_paragraph(); p.space_before = Pt(7); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = name
    r.font.name = MONO if ok is not None else FONT
    r.font.size = Pt(19); r.font.bold = True
    r.font.color.rgb = SOFTWHITE
    p = tf.add_paragraph(); p.space_before = Pt(9); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = what
    r.font.name = FONT; r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0xB6, 0xBC, 0xC4)
    p = tf.add_paragraph(); p.space_before = Pt(11); p.alignment = PP_ALIGN.LEFT
    mark = '' if ok is None else ('\u2713  ' if ok else '\u2717  ')
    r = p.add_run(); r.text = mark + outcome
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = (YELLOW if ok else RGBColor(0xE0, 0x9A, 0x95)) if ok is not None \
        else RGBColor(0xB6, 0xBC, 0xC4)
    cx += 3.05
sh = rect(s, 0.6, 5.55, 12.1, 1.4, RGBColor(0x2A, 0x2D, 0x31))
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.35); tf.margin_right = Inches(0.35)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'Both failures were the **same** failure: I was building the architecture the '
        'corpus knew — an instruction-set simulator written in VHDL — instead of the one '
        'the datasheet drew.', True, 18)
notes(s, 'Now the whole story lands in one place. 297 commits, two complete rewrites. The '
         'diagnosis in the bar is the important sentence and it is not a metaphor: ask a '
         'model to implement an 8008 and everything it has ever seen is an emulator — '
         'sequential software, one opcode at a time. So its instinct pulls that way, and '
         'for two versions I followed it. s8008 is literally an instruction-set simulator '
         'expressed in VHDL. What broke the pattern was going back to the primary source '
         'and building the picture Intel drew.')

# ---- Slide 30: Sweep, part one
s = slide()
kicker(s, 'Part 6 — How this got built')
title(s, 'What it did, and what I had to learn')
COL_A, COL_B = 2.95, 7.75
tb, tf = box(s, COL_A, 1.78, 4.5, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'WHAT THE AI DID'
r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = MIDGREY
tb, tf = box(s, COL_B, 1.78, 4.9, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'WHAT I HAD TO LEARN MYSELF'
r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = RGBColor(0x8A, 0x6D, 0x00)
rows1 = [
    ('What an FPGA is',
     'Could explain any of it, instantly, at any depth. LUTs, routing, block RAM, '
     'bitstreams — perfect answers on demand.',
     'That understanding is not explanation. The book made me **build** the model, in '
     'sequence, with exercises. An AI would have let me skip building it — and then I '
     'would have had nothing to check its work against.'),
    ('What HDL says',
     'Drafted most of the VHDL. Reading and critiquing code is a faster way to learn than '
     'writing from scratch — **this is where it genuinely carried weight.**',
     'The concurrency model. Until you have it you cannot tell correct RTL from *plausible* '
     'RTL — and it produces plausible RTL all day long.'),
    ('Inference',
     'Wrote the block-RAM version. Would just as happily have written the async-read '
     'version, and never mentioned the difference.',
     'Which one was right — from the **utilization report**, not from the code. Both '
     'simulate identically. One does not fit on the chip.'),
]
cy = 2.2
for label, did, learned in rows1:
    rect(s, 0.6, cy, 12.1, 1.42, RGBColor(0xFF, 0xFF, 0xFF),
         line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
    rect(s, 0.6, cy, 0.09, 1.42, YELLOW)
    tb, tf = box(s, 0.85, cy + 0.42, 2.0, 0.8)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(14.5); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    tb, tf = box(s, COL_A, cy + 0.13, 4.6, 1.2)
    p = tf.paragraphs[0]
    rich(p, did, False, 12.5, base_color=RGBColor(0x5A, 0x60, 0x68))
    tb, tf = box(s, COL_B, cy + 0.13, 4.95, 1.2)
    p = tf.paragraphs[0]
    rich(p, learned, False, 12.5, base_color=CHARCOAL)
    cy += 1.5
sh = rect(s, 0.6, cy + 0.06, 12.1, 0.72, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'The pattern: it was fastest where I could already check it — and useless where '
        'I could not.', True, 17)
notes(s, 'This is the promise from the signpost being kept. Walk it, do not read it. Row '
         'one is the least flattering to the AI and it is the one that makes the rest '
         'credible: for foundations it is a tempting shortcut that does not work. You can '
         'get a perfect explanation of a LUT and still not understand FPGAs. Row two is '
         'the honest credit — drafting is real leverage when you can read the output. Row '
         'three is the sharpest concrete example in the talk.')

# ---- Slide 31: Sweep, part two
s = slide()
kicker(s, 'Part 6 — How this got built')
title(s, 'What it did, and what I had to learn  (cont.)')
tb, tf = box(s, COL_A, 1.78, 4.5, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'WHAT THE AI DID'
r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = MIDGREY
tb, tf = box(s, COL_B, 1.78, 4.9, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'WHAT I HAD TO LEARN MYSELF'
r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = RGBColor(0x8A, 0x6D, 0x00)
rows2 = [
    ('The toolchain',
     'Assembled the Makefile. Explained the flags.',
     'To read the **reports**. Routing is two-thirds of my critical path — no model told '
     'me that; nextpnr did. The physics is not in the code, so it is not in the model.'),
    ('φ1 / φ2 as clocks',
     'Would write `rising_edge(phi1)` without hesitation. Valid VHDL. Simulates perfectly.',
     'That it creates a clock domain the timing engine cannot analyze. Learned from a '
     'design that would not close.'),
    ('The 0xFF bug',
     'Nothing. It **could not** — the bug lives below the RTL, inside the synthesizer.',
     'Three days of bisection. And that when simulation and silicon disagree, the '
     'toolchain itself is a suspect.'),
    ('The architecture',
     'Pulled toward an instruction-set simulator every time, because that is all its '
     'corpus contained.',
     '**The block diagram, from the 1972 datasheet.** This is the one that mattered — and '
     'the one thing no model could have handed me.'),
    ('Verification',
     'Wrote much of the 29 testbenches — the boring scaffolding, tirelessly.',
     '`isa.json` and the oracle emulator, by hand from the datasheet. The spec has to come '
     'from the source, not the model — otherwise you are grading its homework with its '
     'own answer key.'),
]
cy = 2.18
for label, did, learned in rows2:
    rect(s, 0.6, cy, 12.1, 0.82, RGBColor(0xFF, 0xFF, 0xFF),
         line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
    rect(s, 0.6, cy, 0.09, 0.82, YELLOW)
    tb, tf = box(s, 0.85, cy + 0.16, 2.0, 0.6)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(13.5); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    tb, tf = box(s, COL_A, cy + 0.06, 4.6, 0.8)
    p = tf.paragraphs[0]
    rich(p, did, False, 11.5, base_color=RGBColor(0x5A, 0x60, 0x68))
    tb, tf = box(s, COL_B, cy + 0.06, 4.95, 0.8)
    p = tf.paragraphs[0]
    rich(p, learned, False, 11.5, base_color=CHARCOAL)
    cy += 0.88
sh = rect(s, 0.6, cy + 0.05, 12.1, 0.7, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'Version three worked because by then I had learned **the right-hand column.**',
     True, 18)
notes(s, 'The payoff. The bar at the bottom is the answer to the question posed back at '
         'the signpost, and it should be said slowly. Row four is the keystone — it ties '
         'to the block-diagram slide and it is the reason the third attempt is different '
         'in kind, not just in quality. Row five is the sentence people write down: '
         'grading its homework with its own answer key. If you are short on time, cut '
         'rows one and two here; four and five carry the argument.')

# ---- Slide 32: What I'd tell you
s = slide()
kicker(s, 'Close')
title(s, 'What I’d tell you about learning this way')
takeaways = [
    ('Read the book first.',
     'Foundations are not delegable. An AI will explain a LUT perfectly and leave you '
     'unable to check its work.'),
    ('Own the specification.',
     'I transcribed every opcode’s T-state sequence from the 1972 manual by hand. The '
     'spec has to come from the primary source.'),
    ('Derive the architecture yourself.',
     'Its corpus had only emulators. The block-based design came from the datasheet — and '
     'it is the thing that made version three work.'),
    ('Keep modules small enough to fully review.',
     '50 to 100 lines, one job. Small enough for me to check, and small enough for it to '
     'get right.'),
    ('Let it draft. You reject.',
     'It made two complete rewrites **affordable** — without that I would have shipped '
     'version one’s architecture and lived with it.'),
    ('Nothing is done until the board says so.',
     'Simulation passing is a hypothesis. That goes double for code you did not write.'),
]
cy = 1.82
for head, body in takeaways:
    rect(s, 0.6, cy, 12.1, 0.76, RGBColor(0xFF, 0xFF, 0xFF),
         line_color=RGBColor(0xD9, 0xDD, 0xE3), line_w=1)
    rect(s, 0.6, cy, 0.09, 0.76, YELLOW)
    tb, tf = box(s, 0.9, cy - 0.01, 11.6, 0.78)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = head
    r.font.name = FONT; r.font.size = Pt(15.5); r.font.bold = True
    r.font.color.rgb = CHARCOAL
    p = tf.add_paragraph()
    rich(p, body, False, 12, base_color=RGBColor(0x5A, 0x60, 0x68))
    cy += 0.83
sh = rect(s, 0.6, cy + 0.02, 12.1, 0.62, CHARCOAL)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'And: 1972 engineers did all of this with 3,500 transistors, no simulator, and '
        'one shot at the mask. Respect.', True, 16)
notes(s, 'Close on the last line — it lands, and it is true. Everything that made this '
         'project possible (simulation, regression tests, a rewritable chip, undo, and a '
         'tireless collaborator) is something the original team did not have. The six '
         'above are the transferable answer to "can you learn a hard skill this way": yes, '
         'and the conditions are specific and non-optional.')

# ---- Slide 28: Q&A
s = slide(dark=True)
title(s, 'Questions?', dark=True, size=44, y=1.9)
logo(s, ASSETS + '/lattice_logo_ondark.png', h=0.34)
bullets(s, 0.6, 3.15, 7.2, 3.3, [
    'Robert Rico',
    'Everything is open: **VHDL, testbenches, the Makefile, the bug reports.**',
    'Start here: **“Getting Started with FPGAs”** — Russell Merrick, No Starch Press',
    'Toolchain: GHDL · Yosys · nextpnr · Project Trellis · openFPGALoader',
    'Period software: Mike Willegal’s SCELBI archive · SCELBAL port by Jim Loos',
], dark=True, size=15.5, space=13)
sh = rect(s, 0.6, 6.5, 7.2, 0.8, CHARCOAL2)
tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
rich(p, 'The board stays up afterwards — come type at it.', True, 16)
if have(REPO + '/docs/images/tiny_os_scelbal.png'):
    tw, th = fit(REPO + '/docs/images/tiny_os_scelbal.png', 4.7, 3.4)
    pic(s, REPO + '/docs/images/tiny_os_scelbal.png', 8.2, 2.6, w=tw)
notes(s, 'Fill in the repo URL when public. Likely questions: "why not Verilog/SystemVerilog" '
         '(taste + type strictness), "why not a Xilinx board" (open toolchain), '
         '"how much did the AI actually write" (most of the VHDL drafts; all of the '
         'specification, verification, and rejection was mine), "could you run it faster" '
         '(yes — 190× headroom — but then the period software’s timing loops break).')

prs.save(OUT)
print(f'Wrote {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides')
