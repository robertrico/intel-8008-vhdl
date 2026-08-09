#!/usr/bin/env python3
"""Build b8008 presentation deck per docs/presentation/deck_spec.md

Two-zone theme:
  Intel zone  (slides 1-8):  Intel Blue structure, deep Intel-blue dark slides,
                             energy-blue accents. No yellow anywhere.
  Lattice zone (slides 9+):  charcoal + yellow. Theme flips on the FPGA primer
                             slide (dark, Lattice logo) - "entering Lattice world".
"""
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = '/Users/hackbook/Development/intel-8008-vhdl'
ASSETS = REPO + '/docs/presentation/assets'
OUT = REPO + '/docs/presentation/b8008_talk.pptx'

INTEL = RGBColor(0x00, 0x68, 0xB5)
INTEL_DARK = RGBColor(0x00, 0x42, 0x77)   # deep Intel blue, dark-slide bg
ENERGY = RGBColor(0x00, 0xC7, 0xFD)       # Intel 2020 "energy blue" accent
YELLOW = RGBColor(0xFF, 0xC2, 0x20)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
CHARCOAL2 = RGBColor(0x3B, 0x3F, 0x44)
LIGHTBG = RGBColor(0xF7, 0xF8, 0xFA)
SOFTWHITE = RGBColor(0xF2, 0xF2, 0xF2)
CODEBG = RGBColor(0x2A, 0x2D, 0x31)
SKY = RGBColor(0x1F, 0x9D, 0xD8)
MIDGREY = RGBColor(0x8D, 0x8D, 0x8D)
CODE_COMMENT = RGBColor(0x6F, 0xC2, 0xEA)
PALE_BLUE = RGBColor(0xCD, 0xEB, 0xFA)
PALE_YELLOW = RGBColor(0xFF, 0xF3, 0xD3)
GREYFILL = RGBColor(0xE9, 0xEC, 0xF0)

FONT = 'Aptos'
MONO = 'Aptos Mono'

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

ZONE = 'intel'   # flips to 'lattice' at the FPGA primer


def z_dark_bg():
    return INTEL_DARK if ZONE == 'intel' else CHARCOAL


def z_primary():          # structural color on light slides
    return INTEL if ZONE == 'intel' else CHARCOAL


def z_accent_dark():      # accent on dark slides
    return ENERGY if ZONE == 'intel' else YELLOW


def z_accent_light():     # accent bar color on light slides
    return INTEL if ZONE == 'intel' else YELLOW


def z_code_inline():      # inline `code` color on light slides
    return INTEL if ZONE == 'intel' else SKY


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = z_dark_bg() if dark else LIGHTBG
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def style_run(r, dark, size, mono=False, bold=False, color=None):
    r.font.name = MONO if mono else FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color if color else (SOFTWHITE if dark else CHARCOAL)


def rich(p, text, dark, size, base_color=None):
    """Parse **bold** and `mono` runs into paragraph p."""
    for tok in re.split(r'(\*\*.*?\*\*|`.*?`)', text):
        if not tok:
            continue
        if tok.startswith('**') and tok.endswith('**'):
            r = p.add_run(); r.text = tok[2:-2]
            style_run(r, dark, size, bold=True, color=base_color)
        elif tok.startswith('`') and tok.endswith('`'):
            r = p.add_run(); r.text = tok[1:-1]
            style_run(r, dark, size, mono=True,
                      color=z_accent_dark() if dark else z_code_inline())
        else:
            r = p.add_run(); r.text = tok
            style_run(r, dark, size, color=base_color)


def bullets(s, x, y, w, h, lines, dark=False, size=18, space=8, anchor_top=True):
    tb, tf = box(s, x, y, w, h)
    if not anchor_top:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space)
        indent = line.startswith('  ')
        txt = line.strip()
        if txt.startswith('- '):
            txt = ('–  ' if indent else '•  ') + txt[2:]
            if indent:
                p.level = 1
        rich(p, txt, dark, size - (2 if indent else 0))
    return tb


def kicker(s, text, dark=False):
    tb, tf = box(s, 0.55, 0.32, 8.0, 0.35)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text.upper()
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
    if dark:
        r.font.color.rgb = z_accent_dark()
    else:
        r.font.color.rgb = INTEL if ZONE == 'intel' else CHARCOAL2


def title(s, text, dark=False, size=32, y=0.62, w=12.2):
    tb, tf = box(s, 0.55, y, w, 1.0)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = SOFTWHITE if dark else z_primary()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y + 0.78),
                             Inches(1.6), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = z_accent_dark() if dark else z_accent_light()
    bar.line.fill.background()
    return tb


def rect(s, x, y, w, h, fill, line_color=None, line_w=None,
         shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w or 1)
    sh.shadow.inherit = False
    return sh


def shape_text(sh, text, size=14, bold=False, color=None, mono=False,
               align=PP_ALIGN.CENTER):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = MONO if mono else FONT
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color if color else SOFTWHITE


def placeholder(s, x, y, w, h, label, dark=False):
    sh = rect(s, x, y, w, h, CHARCOAL2 if dark else RGBColor(0xEA, 0xED, 0xF1),
              line_color=MIDGREY, line_w=1.25)
    from pptx.oxml.ns import qn
    ln = sh.line._get_or_add_ln()
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(d)
    shape_text(sh, 'PLACEHOLDER\n' + label, size=13, bold=False,
               color=SOFTWHITE if dark else RGBColor(0x6A, 0x6F, 0x76))
    return sh


def code_panel(s, x, y, w, h, code, size=14, bg=None):
    sh = rect(s, x, y, w, h, bg or CODEBG)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.12)
    for i, ln in enumerate(code.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if '--' in ln:
            main, com = ln.split('--', 1)
        elif ln.strip().startswith('//'):
            main, com = '', ln
        else:
            main, com = ln, None
        if main:
            r = p.add_run(); r.text = main
            r.font.name = MONO; r.font.size = Pt(size)
            r.font.color.rgb = SOFTWHITE
        if com is not None:
            r = p.add_run(); r.text = ('--' + com) if '--' in ln else com
            r.font.name = MONO; r.font.size = Pt(size)
            r.font.color.rgb = CODE_COMMENT
        if main == '' and com is None:
            r = p.add_run(); r.text = ''
            r.font.name = MONO; r.font.size = Pt(size)
    return sh


def pic(s, path, x, y, w=None, h=None, border=True):
    kw = {}
    if w: kw['width'] = Inches(w)
    if h: kw['height'] = Inches(h)
    p = s.shapes.add_picture(path, Inches(x), Inches(y), **kw)
    if border:
        p.line.color.rgb = MIDGREY
        p.line.width = Pt(0.75)
    p.shadow.inherit = False
    return p


def img_size(path):
    return Image.open(path).size


def fit(path, max_w, max_h):
    w, h = img_size(path)
    scale = min(max_w / w, max_h / h)
    return w * scale, h * scale


def logo(s, path, h=0.4, x=None, y=0.35):
    """Corner brand logo, top-right by default."""
    iw, ih = img_size(path)
    w = h * iw / ih
    if x is None:
        x = 13.333 - w - 0.55
    return pic(s, path, x, y, h=h, border=False)


def arrow(s, x1, y1, x2, y2, color=None, wpt=2.25):
    conn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color if color else z_primary()
    conn.line.width = Pt(wpt)
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    head = ln.makeelement(qn('a:tailEnd'), {'type': 'arrow'})
    ln.append(head)
    conn.shadow.inherit = False
    return conn


def table(s, x, y, w, h, rows, col_widths, header=True, size=16):
    shp = s.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y),
                             Inches(w), Inches(h))
    t = shp.table
    hdr_color = INTEL if ZONE == 'intel' else CHARCOAL2
    alt = RGBColor(0xED, 0xF1, 0xF6) if ZONE == 'intel' else RGBColor(0xEF, 0xEF, 0xEF)
    for ci, cw in enumerate(col_widths):
        t.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            if header and ri == 0:
                cell.fill.fore_color.rgb = hdr_color
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 else alt
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            rich(p, str(val), False, size,
                 base_color=RGBColor(0xFF, 0xFF, 0xFF) if (header and ri == 0) else CHARCOAL)
            if header and ri == 0:
                for r in p.runs:
                    r.font.bold = True
    return t


# ============================================================
# INTEL ZONE — slides 1-8
# ============================================================

# ---- Slide 1: Title (Intel dark blue, C8008 chip photo)
s = slide(dark=True)
tw, th = fit(ASSETS + '/intel_c8008_chip.jpg', 5.6, 3.9)
pic(s, ASSETS + '/intel_c8008_chip.jpg', 13.333 - tw - 0.7, (7.5 - th) / 2, w=tw)
tb, tf = box(s, 0.7, 2.2, 6.4, 3.4)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'b8008'
r.font.name = MONO; r.font.size = Pt(54); r.font.bold = True
r.font.color.rgb = ENERGY
p = tf.add_paragraph(); p.space_before = Pt(6)
r = p.add_run(); r.text = 'Resurrecting the World’s First\n8-Bit Microprocessor'
r.font.name = FONT; r.font.size = Pt(30); r.font.bold = True
r.font.color.rgb = SOFTWHITE
p = tf.add_paragraph(); p.space_before = Pt(18)
r = p.add_run(); r.text = 'A cycle-exact Intel 8008 in VHDL,\nbuilt with an AI pair programmer'
r.font.name = FONT; r.font.size = Pt(17)
r.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE2)
p = tf.add_paragraph(); p.space_before = Pt(22)
r = p.add_run(); r.text = 'Robert Rico'
r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
r.font.color.rgb = SOFTWHITE
rect(s, 0.75, 2.15, 1.6, 0.045, ENERGY)
# brand row: 1968 dropped-e Intel logo (period-correct for a 1972 chip) +
# Lattice, each underlined in its zone color — silent legend for the theme
lh = 0.5
iw2, ih2 = img_size(ASSETS + '/intel_1968_ondark.png')
i_w = lh * iw2 / ih2
lw2, lh2 = img_size(ASSETS + '/lattice_logo_ondark.png')
l_h = 0.32
l_w = l_h * lw2 / lh2
ly = 6.35
pic(s, ASSETS + '/intel_1968_ondark.png', 0.78, ly, h=lh, border=False)
rect(s, 0.78, ly + lh + 0.12, i_w, 0.045, ENERGY)
lx = 0.78 + i_w + 0.75
pic(s, ASSETS + '/lattice_logo_ondark.png', lx, ly + (lh - l_h) / 2, h=l_h, border=False)
rect(s, lx, ly + lh + 0.12, l_w, 0.045, YELLOW)
notes(s, 'One-liner intro. "In 1972 Intel shipped the first 8-bit microprocessor. '
         'I spent 8 months bringing it back — and tonight it’s running on the table." '
         'Deck color story: Intel blue until we reach the FPGA; charcoal+yellow (Lattice) '
         'from there — the theme itself says whose world we’re in. Never say it out loud. '
         'Chip photo: Konstantin Lanzet, Wikimedia Commons, CC BY-SA 3.0.')

# ---- Slide 2: What is the 8008 (Intel logo re-introduced)
s = slide()
kicker(s, 'Act I — What')
title(s, 'The Intel 8008 (April 1972)')
logo(s, ASSETS + '/intel_1968_onlight.png', h=0.52)
bullets(s, 0.6, 1.9, 7.4, 5.0, [
    '- World’s first 8-bit microprocessor',
    '- 3,500 transistors, 10 µm PMOS',
    '- **18 pins** — for a CPU',
    '- 500 kHz clock; instructions take 5, 8, or 11 T-states',
    '- 16 KB address space (14-bit)',
    '- Seven 8-bit registers (A, B, C, D, E, H, L)',
    '- 48 instructions',
], size=19, space=10)
cw2, ch2 = fit(ASSETS + '/intel_c8008_chip.jpg', 4.3, 4.6)
pic(s, ASSETS + '/intel_c8008_chip.jpg', 8.4, 1.9 + (4.6 - ch2) / 2, w=cw2)
notes(s, 'Emphasize 18 pins — that single number causes most of the weirdness coming up. '
         'Modern context: an ATtiny has more pins. '
         'Chip photo: Konstantin Lanzet, Wikimedia Commons, CC BY-SA 3.0.')

# ---- Slide 3: History
s = slide()
kicker(s, 'Act I — What')
title(s, 'Built for a terminal, rejected by its buyer')
bullets(s, 0.6, 1.9, 7.6, 5.2, [
    '- 1969-70: **CTC (Computer Terminal Corp., later Datapoint)** designs the Datapoint 2200 "programmable terminal" — CPU built from ~100 TTL chips',
    '- CTC asks **Intel** and **Texas Instruments** to shrink it to one chip',
    '- TI’s **TMX 1795** (1971): technically first 8-bit CPU chip — never shipped',
    '- Intel’s chip: late, and slower than CTC’s TTL board — **CTC walks away**',
    '- Intel keeps the rights in trade, sells it as the **8008**',
], size=18, space=12)
placeholder(s, 8.5, 1.9, 4.2, 4.6, 'Datapoint 2200 photo\n(Wikimedia Commons)\n\noptional: TMX 1795 die\n(Ken Shirriff)')
notes(s, 'The ISA was CTC’s design, not Intel’s. TI beat everyone to silicon but never '
         'productized — the 8008 outlasted it by shipping.')

# ---- Slide 4: Lineage
s = slide()
kicker(s, 'Act I — What')
title(s, 'Your laptop still carries its DNA')
chips = [('8008', '1972'), ('8080', '1974'), ('8086', '1978'), ('x86', 'today')]
cx = 0.9
for i, (name, yr) in enumerate(chips):
    last = i == len(chips) - 1
    sh = rect(s, cx, 2.15, 2.15, 1.15, INTEL_DARK if last else INTEL,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(sh, f'{name}\n{yr}', size=18, bold=True,
               color=ENERGY if last else SOFTWHITE)
    if not last:
        arrow(s, cx + 2.15 + 0.05, 2.72, cx + 2.15 + 0.85, 2.72)
    cx += 3.0
bullets(s, 0.9, 3.9, 11.5, 1.6, [
    'The instruction-set lineage in every PC today traces back to a terminal company in San Antonio, Texas.',
    'And in 1974, the **SCELBI-8H** — arguably the first personal computer — was built around this chip. That’s where our demo software comes from.',
], size=18, space=12)
placeholder(s, 9.1, 5.35, 3.6, 1.75, 'optional: SCELBI-8H photo (willegal.net)')
notes(s, 'Bridge slide — plants SCELBI now so the demos later have context.')

# ---- Slide 5: Weirdness 1 clocks
s = slide()
kicker(s, 'Act I — The Quirks')
title(s, 'Weirdness #1 — Two clocks, never high together')
bullets(s, 0.6, 1.95, 5.6, 4.4, [
    'The 8008 needs **two non-overlapping clock phases** (φ1, φ2).',
    'No single clock pin. External circuitry must generate both phases with guaranteed dead time between them.',
    'Get the overlap wrong → internal buses fight → **undefined behavior**.',
], size=18, space=14)
ow, oh = fit(REPO + '/docs/two_phase_oscope_cap_1.png', 6.3, 2.45)
pic(s, REPO + '/docs/two_phase_oscope_cap_1.png', 6.5, 1.9, w=ow)
pic(s, REPO + '/docs/two_phase_oscope_cap_2.png', 6.5, 1.9 + oh + 0.22, w=ow)
tb, tf = box(s, 6.5, 1.9 + 2 * oh + 0.3, 6.3, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'Real capture: the FPGA generating φ1/φ2 on my bench'
r.font.name = FONT; r.font.size = Pt(12); r.font.italic = True
r.font.color.rgb = MIDGREY
notes(s, 'These scope shots are from my bench — this is the FPGA generating the two phases. '
         'In 1972 you built this from discrete logic before the CPU would even blink.')

# ---- Slide 6: Weirdness 2 bus
s = slide()
kicker(s, 'Act I — The Quirks')
title(s, 'Weirdness #2 — One bus, three jobs')
bullets(s, 0.6, 1.85, 5.9, 1.4, [
    '18 pins means **8 data pins total**. So the bus is time-multiplexed, three ways, every machine cycle:',
], size=18)
table(s, 0.6, 3.1, 5.9, 2.0, [
    ['T-state', 'What’s on the 8 pins'],
    ['T1', 'Address, low byte'],
    ['T2', 'Address high 6 bits + cycle type'],
    ['T3', 'The actual data'],
], [1.4, 4.5], size=15)
bullets(s, 0.6, 5.55, 5.9, 1.0, [
    'External latches must grab each piece at the right instant.',
], size=18)
tw2, th2 = fit(ASSETS + '/8008_instruction_cycle.png', 6.1, 4.7)
pic(s, ASSETS + '/8008_instruction_cycle.png', 6.8, 1.85, w=tw2)
tb, tf = box(s, 6.8, 1.85 + th2 + 0.08, 6.1, 0.35)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'Fig. 1, 8008 User’s Manual — one instruction cycle'
r.font.name = FONT; r.font.size = Pt(12); r.font.italic = True
r.font.color.rgb = MIDGREY
notes(s, '"Tri-plexed" bus. Every memory access is a little three-act play. The CPU also '
         'broadcasts its internal state on pins S0-S2 so external logic can follow along.')

# ---- Slide 7: Weirdness 3 boot by interrupt
s = slide()
kicker(s, 'Act I — The Quirks')
title(s, 'Weirdness #3 — It boots by being interrupted')
bullets(s, 0.6, 1.9, 5.7, 4.8, [
    '- No reset vector. At power-on the PC is garbage.',
    '- The *only* way to start it: assert the interrupt line and **jam an instruction onto the bus** (typically `RST 0`)',
    '- Startup is an interrupt. Interrupts are instruction injection. It’s injection all the way down.',
    '- Intel’s own reference design (SIM8-01) needed **dozens of TTL support chips** just to make the CPU usable.',
], size=17, space=13)
sw2, sh2 = fit(ASSETS + '/sim8_01_schematic.png', 6.4, 4.4)
pic(s, ASSETS + '/sim8_01_schematic.png', 6.5, 1.9, w=sw2)
tb, tf = box(s, 6.5, 1.9 + sh2 + 0.08, 6.4, 0.4)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'SIM8-01: find the CPU. (Hint: the small box in the middle.)'
r.font.name = FONT; r.font.size = Pt(12); r.font.italic = True
r.font.color.rgb = MIDGREY
notes(s, 'This is why "just wire up the CPU" was never an option — in 1972 or now. '
         'The support glue is half the project.')

# ---- Slide 8: Weirdness 4 PC-in-stack
s = slide()
kicker(s, 'Act I — The Quirks')
title(s, 'Weirdness #4 — There is no program counter')
bullets(s, 0.6, 1.9, 6.6, 5.0, [
    '…there are **eight** of them.',
    'The PC is whichever slot of an 8 × 14-bit address stack the stack pointer selects.',
    '- **CALL** = SP moves forward (old slot keeps the return address)',
    '- **RET** = SP moves back',
    '- 7 levels of nesting; the 8th CALL silently wraps onto the oldest',
    'No stack in memory. No PUSH/POP of data. **1972 had opinions.**',
], size=17, space=11)
sx, sy, sw3, rh = 8.6, 1.75, 3.0, 0.58
for i in range(8):
    active = i == 2
    sh = rect(s, sx, sy + i * (rh + 0.06), sw3, rh,
              PALE_BLUE if active else GREYFILL,
              line_color=INTEL if active else MIDGREY, line_w=1.5 if active else 1)
    lbl = f'slot {7 - i}'
    if active:
        lbl += '   ← SP  (this IS the PC)'
    shape_text(sh, lbl, size=13, bold=active, color=INTEL if active else CHARCOAL)
tb, tf = box(s, sx - 0.1, sy + 8 * (rh + 0.06) + 0.05, 3.6, 0.6)
p = tf.paragraphs[0]
r = p.add_run(); r.text = '8 × 14-bit address stack'
r.font.name = FONT; r.font.size = Pt(12); r.font.italic = True
r.font.color.rgb = MIDGREY
notes(s, 'Remember this one — it becomes the heart of the final design later. Implemented '
         'structurally, exactly as drawn, the stack-wrap behavior falls out for free.')

# ============================================================
# LATTICE ZONE — slide 9 onward (theme flip)
# ============================================================
ZONE = 'lattice'

# ---- Slide 9: The FPGA I chose (Versa board) — entering Lattice world
s = slide(dark=True)
kicker(s, 'Act II — FPGA Primer', dark=True)
title(s, 'The FPGA I chose: Lattice ECP5-5G Versa', dark=True)
logo(s, ASSETS + '/lattice_logo_ondark.png', h=0.34)
bullets(s, 0.6, 1.9, 6.9, 4.3, [
    '- **LFE5UM5G-45F** — 44K LUTs, 1.9 Mb block RAM, 5 Gbps SerDes',
    '- Dev board: DIP switches + LEDs become the 8008’s **front panel**; UART over USB is the terminal',
    '- Best-in-class **open-source toolchain** support: Yosys + nextpnr know this chip inside out',
    '- Design runs at 25 MHz system clock; timing closes at ~95 MHz. The 8008 needs **500 kHz**.',
], dark=True, size=17, space=14)
sh = rect(s, 0.6, 6.1, 6.9, 0.85, CHARCOAL2)
tf2 = sh.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'The entire 1972 computer system: 1,672 LUTs — '
r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
r.font.color.rgb = SOFTWHITE
r = p.add_run(); r.text = '3% of the chip'
r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True
r.font.color.rgb = YELLOW
placeholder(s, 7.9, 1.9, 4.8, 4.6, 'ECP5-5G Versa board photo\nor product shot / logo\n(user supplies)', dark=True)
notes(s, 'THEME FLIP: we just left Intel’s world (blue) and entered Lattice’s (charcoal + '
         'yellow) — the deck stays this way to the end. Numbers from the real b8008_basic '
         'build reports: 1,672/43,848 LUT4s (3%), 8/108 block RAMs (7%), fmax ~95 MHz vs '
         '25 MHz system clock. The chip is comically oversized for the job — the whole 1972 '
         'rack fits in a corner of it.')

# ---- Slide 10: FPGA primer
s = slide(dark=True)
kicker(s, 'Act II — FPGA Primer', dark=True)
title(s, '30-second FPGA primer', dark=True)
bullets(s, 0.6, 1.9, 7.3, 4.8, [
    '- A chip full of uncommitted logic blocks and programmable wiring',
    '- You describe a circuit; the toolchain configures the chip to *become* it',
    '- Not emulation, not software — it **is** the circuit, in real silicon, in parallel',
    '- This project targets a **Lattice ECP5-5G** using a fully open-source toolchain (GHDL + Yosys + nextpnr)',
], dark=True, size=18, space=14)
placeholder(s, 8.3, 1.9, 4.4, 4.6, 'FPGA fabric diagram\n(generic LUT/routing illustration)\nor draw grid graphic', dark=True)
notes(s, 'For the software folks: the 8008 design isn’t simulated on the board — the board '
         'is wired up as an 8008. All that 1972 support glue — clocks, latches, ROM, RAM, '
         'UART — fits inside the same chip.')

# ---- Slide 11: VHDL vs Verilog
s = slide()
kicker(s, 'Act II — FPGA Primer')
title(s, 'Hardware description languages in one slide')
bullets(s, 0.6, 1.8, 12.1, 1.15, [
    'Two mainstream ways to describe digital hardware — same job: describe registers and logic; tools synthesize real gates.',
], size=18)
tb, tf = box(s, 0.7, 2.9, 5.8, 0.4)
p = tf.paragraphs[0]; rich(p, '**Verilog** — C-flavored, terse, dominant in US industry', False, 16)
code_panel(s, 0.7, 3.35, 5.8, 2.3,
"""always @(posedge clk)
  if (en)
    count <= count + 1;""", size=15)
tb, tf = box(s, 6.9, 2.9, 5.8, 0.4)
p = tf.paragraphs[0]; rich(p, '**VHDL** — Ada-flavored, verbose, strongly typed', False, 16)
code_panel(s, 6.9, 3.35, 5.8, 2.9,
"""process(clk) begin
  if rising_edge(clk) then
    if en = '1' then
      count <= count + 1;
    end if;
  end if;
end process;""", size=15)
bullets(s, 0.7, 6.5, 12.0, 0.7, [
    'This project: **VHDL**, one file per functional block.',
], size=18)
notes(s, 'Don’t linger. "Verbose but strict — it catches a class of mistakes at compile '
         'time, which matters when your collaborator is an AI."')

# ---- Slide 11: The project
s = slide()
kicker(s, 'Act III — The AI-Aided Project')
title(s, 'The project: 9 months, 3 attempts')
bullets(s, 0.6, 1.85, 12.2, 1.3, [
    '**Nov 2025 → Aug 2026** — 367 commits. **3 iterations, 2 complete bottom-to-top rewrites.**',
], size=19)
ty = 3.15
attempts = [
    ('s8008', 'monolithic, single-cycle\nWorked in simulation.\nBroke on hardware.', False, 3.4),
    ('v8008', 'multi-cycle rewrite\nCollapsed under its own\ncomplexity.', False, 3.4),
    ('b8008', 'block-based, cycle-exact\nSILICON VALIDATED', True, 4.2),
]
tx = 0.8
for name, desc, ok, w in attempts:
    sh = rect(s, tx, ty, w, 1.85, CHARCOAL if ok else GREYFILL,
              line_color=None if ok else MIDGREY, line_w=1,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf2 = sh.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ('✓  ' if ok else '✗  ') + name
    r.font.name = MONO; r.font.size = Pt(19); r.font.bold = True
    r.font.color.rgb = YELLOW if ok else CHARCOAL2
    for ln in desc.split('\n'):
        p = tf2.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(12.5)
        r.font.color.rgb = SOFTWHITE if ok else CHARCOAL
    if not ok:
        arrow(s, tx + w + 0.03, ty + 0.92, tx + w + 0.42, ty + 0.92, color=CHARCOAL)
    tx += w + 0.45
bullets(s, 0.6, 5.5, 12.2, 1.4, [
    'Built in collaboration with **Claude** (Anthropic’s AI) — every line reviewed, tested, and frequently rejected',
], size=19)
notes(s, 'Set expectations honestly: this was NOT "AI, build me an 8008." Nine months of '
         'iteration, and I threw the whole thing away twice. Here’s what the first two '
         'attempts taught me.')

# ---- Slide 12: s8008 lesson
s = slide()
kicker(s, 'Act III — The AI-Aided Project')
title(s, 's8008: when "all tests pass" lies to you')
bullets(s, 0.6, 1.95, 11.9, 4.6, [
    'The monolithic version **passed its whole simulation suite**.',
    'On the FPGA, ALU operations returned garbage.',
    '- Single-cycle timing masked bugs simulation couldn’t see',
    '- One 1,500-line file nobody — human or AI — could fully reason about',
], size=19, space=13)
sh = rect(s, 0.6, 5.3, 11.9, 1.0, PALE_YELLOW)
shape_text(sh, 'Rule adopted: nothing is "done" until it’s proven on silicon.',
           size=20, bold=True, color=CHARCOAL)
notes(s, 'This failure defined the whole methodology afterward: sim-pass is a hypothesis, '
         'hardware is the proof. The v8008 rewrite then failed differently — too clever, too '
         'entangled. Two failures, one diagnosis: complexity nobody could reason about.')

# ---- Slide 13: Working with AI
s = slide()
kicker(s, 'Act III — The AI-Aided Project')
title(s, 'The AI wrote most of the VHDL. Here’s why that worked.')
bullets(s, 0.6, 1.9, 6.9, 4.6, [
    'It worked because **I could tell when it was wrong**.',
    '- I had to learn the 8008 cold: T-states, machine cycles, every ISA quirk',
    '- Built `isa.json` — every opcode mapped to its documented T-state sequence — as ground truth',
    '- 37-test regression suite; cycle counts diffed against the datasheet automatically',
    '**The AI multiplied effort. The verification was the job.**',
], size=16.5, space=11)
code_panel(s, 7.8, 1.9, 4.9, 4.9,
"""// docs/isa.json - ground truth
{ "operation": "CAL",
  "D[7:0]": "01XXX110",
  "num_states_execution": 11,
  "num_cycles": 3,
  "cycle_1": {
    "T1": "PCL OUT",
    "T2": "PCH OUT",
    "T3": "FETCH INSTR.",
    ...
  "cycle_3": {
    ...
    "T4": "REG.a TO PCH",
    "T5": "REG.b TO PCL" }}""", size=12.5)
notes(s, 'The transferable lesson for this audience: AI collaboration shifts your work from '
         'writing to specifying and verifying. You must know enough to catch it being '
         'confidently wrong — and it will be. SEGUE: "So attempt #3 needed a design that a '
         'human AND an AI could each fully reason about. Turns out Intel drew it for us in 1972."')

# ---- Slide 14: Block diagram
s = slide()
kicker(s, 'Act IV — The Design the Failures Forced')
title(s, 'The fix: build Intel’s own block diagram. Literally.')
bw, bh = fit(ASSETS + '/8008_block_diagram.png', 9.4, 5.0)
pic(s, ASSETS + '/8008_block_diagram.png', 0.6, 1.85, w=bw)
bullets(s, 0.6 + bw + 0.35, 2.3, 12.7 - bw - 0.95, 4.0, [
    'Attempt #3:',
    '**one module per block.**',
    '**No cleverness.**',
    'Fig. 3, 8008 User’s Manual (1973)',
], size=17, space=12)
notes(s, 'After two failures from complexity, the answer was radical simplicity: stop inventing '
         'an architecture and copy the one on the datasheet. Walk the diagram: address stack, '
         'scratchpad registers, ALU, instruction register/decoder, timing & control. Every box '
         'becomes a VHDL file.')

# ---- Slide 15: Dumb modules
s = slide()
kicker(s, 'Act IV — The Design the Failures Forced')
title(s, 'Design rule: every module is dumb')
bullets(s, 0.6, 1.9, 6.2, 4.8, [
    '**25 VHDL files.** Each block:',
    '- Does **one job**, ~50–100 lines',
    '- Knows **nothing** about instructions, interrupts, or other modules',
    '- Responds only to explicit control signals',
    '- Has its own testbench',
    'All the intelligence lives in one place: the **timing & control unit**.',
], size=17, space=12)
ctl = rect(s, 7.3, 1.85, 5.3, 0.85, CHARCOAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(ctl, 'Timing & Control Unit\n(the only smart module)', size=13, bold=True,
           color=YELLOW)
mods = [('address\nstack', 0), ('stack\npointer', 1), ('register\nfile', 2),
        ('ALU', 3), ('instr.\ndecoder', 4)]
mw = 0.97
for name, i in mods:
    mx = 7.3 + i * (mw + 0.115)
    arrow(s, mx + mw / 2, 2.75, mx + mw / 2, 3.35, color=CHARCOAL, wpt=1.5)
    sh = rect(s, mx, 3.4, mw, 0.95, GREYFILL, line_color=MIDGREY, line_w=1)
    shape_text(sh, name, size=11, bold=True, color=CHARCOAL)
tb, tf = box(s, 7.3, 4.5, 5.3, 0.9)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'control signals only — no module knows what a "JMP" is'
r.font.name = FONT; r.font.size = Pt(12); r.font.italic = True
r.font.color.rgb = MIDGREY
sh = rect(s, 7.3, 5.5, 5.3, 1.0, PALE_YELLOW)
shape_text(sh, 'Mirrors Fig. 3 box-for-box —\nthe diagram IS the architecture', size=14,
           bold=True, color=CHARCOAL)
notes(s, 'Direct answer to the s8008 postmortem: when a module can’t know about instructions, '
         'it can’t harbor an instruction-specific bug. Debugging becomes: which control signal '
         'was wrong, and when? Also the AI payoff: a 70-line module with one job is something '
         'Claude gets right — and something I can fully review.')

# ---- Slide 16: Stack pointer code
s = slide()
kicker(s, 'Act IV — The Design the Failures Forced')
title(s, 'The whole stack pointer. All of it.')
code_panel(s, 0.6, 1.85, 7.6, 4.6,
"""if reset = '1' then
    sp <= (others => '0');
elsif rising_edge(clk) then
    if phi1_rising = '1' then
        if stack_push = '1' then
            sp <= sp + 1;   -- wraps 7 -> 0
        elsif stack_pop = '1' then
            sp <= sp - 1;   -- wraps 0 -> 7
        end if;
    end if;
end if;""", size=17)
bullets(s, 8.5, 2.1, 4.2, 4.2, [
    '`src/b8008/stack_pointer.vhdl`',
    '**69 lines** including comments.',
    'It has never had a bug.',
    'It doesn’t know what CALL is. It doesn’t know interrupts exist.',
    '**It counts.**',
], size=16, space=13)
notes(s, 'CALL/RET/RST semantics — even the stack-wrap behavior from Weirdness #4 — emerge '
         'from structure, not from logic.')

# ---- Slide 17: Litmus test
s = slide()
kicker(s, 'Act V — How Do I Know It Works?')
title(s, 'The real test suite was written in 1974')
bullets(s, 0.6, 1.85, 7.5, 1.15, [
    'Directed tests prove what you thought to test.',
    '**Period software proves the machine.**',
], size=19, space=6)
bullets(s, 0.6, 3.15, 7.5, 3.0, [
    '- SCELBI floating-point calculator (1974)',
    '- HEXPAWN — self-modifying learning game (1973)',
    '- STARS (Byte magazine, 1976)',
    '- Pi digit generator, Mandelbrot renderer',
    '- SCELBAL BASIC (1976)',
], size=18, space=8)
bullets(s, 0.6, 6.15, 11.9, 1.0, [
    'Each ported with a minimal, documented change ledger — I/O shims and memory-map constants only.',
], size=16)
placeholder(s, 8.6, 1.85, 4.1, 4.1, 'SCELBI book cover or\nByte May 1976 cover\n(period flavor)')
notes(s, 'If a self-modifying 1973 game and a floating-point interpreter both run, the CPU is '
         'right in ways no directed test can claim. This software exercised corners I’d never '
         'have thought to test.')

# ---- Slide 18: Scorecard
s = slide()
kicker(s, 'Act V — How Do I Know It Works?')
title(s, 'Verification scorecard')
table(s, 1.2, 2.0, 10.9, 4.4, [
    ['Check', 'Result'],
    ['Regression suite', '37/37'],
    ['Cycle-exact T-states vs datasheet (all 27 timing classes)', '27/27'],
    ['Interrupt suite', '10/10'],
    ['Hardware ISA self-test, running on the board', '46/46'],
    ['Flags spec-exact (INR/DCR preserve carry; rotates carry-only)', '✓'],
    ['Interrupts at instruction boundaries only (User’s Manual Fig. 2)', '✓'],
], [8.4, 2.5], size=16)
notes(s, '"Cycle-exact" means: not just the right answer — the right answer in exactly the '
         'number of clock states the 1972 datasheet specifies, for every instruction class. '
         'Optional: swap table for green run_all_tests.sh terminal screenshot.')

# ---- Slide 19b: Verification tooling (formal proofs, equivalence, fuzzing)
s = slide()
kicker(s, 'Act V — How Do I Know It Works?')
title(s, 'Beyond tests: proofs, equivalence, fuzzing')
table(s, 0.6, 1.85, 12.1, 3.9, [
    ['Layer', 'What it checks', 'Scale'],
    ['SBY property proofs', 'Module contracts as PSL assertions, proven by SMT solver '
     '(k-induction / bounded)', '11 suites'],
    ['Equivalence checks', 'RTL vs its synthesized gate netlist (Yosys round trip)',
     '7 miters + 6 EQY'],
    ['Exhaustive sweeps', 'ALU and instruction decoder vs independent Python models',
     '1,049,600 + 256 cases'],
    ['cocotb monitors', 'External bus protocol on the full core; per-instruction '
     'control-signal scenarios', 'RTL + netlist'],
    ['Differential fuzzer', 'Random legal programs under three oracles: bus monitor, '
     'datasheet timing, RTL-vs-netlist trace diff', 'seeded, both cores'],
], [2.5, 7.2, 2.4], size=13.5)
sh = rect(s, 0.6, 6.0, 12.1, 0.85, PALE_YELLOW)
shape_text(sh, '37 CI jobs on every push  ·  every checker mutation-tested  ·  '
              'verification plan: 102 rows, zero gaps', size=16, bold=True, color=CHARCOAL)
notes(s, 'One paragraph per row, no deep dive. SBY: the solver proves the assertion for ALL '
         'input sequences, not sampled ones — math, not test vectors. Equivalence: the exact '
         'netlist headed for the FPGA is proven to match the RTL, so synthesis can’t silently '
         'change behavior. Exhaustive: every ALU input combination, every opcode, against '
         'models written from the datasheet, not from the VHDL. Mutation-tested = plant a bug, '
         'watch the checker fail, revert — a checker that can’t fail proves nothing. '
         'Anecdote: the bus monitor found a real RTL bug — two machine-cycle type codes '
         'transposed on the external bus. The fuzzer runs the same random program on the RTL '
         'core and the synthesized-netlist core and diffs the traces.')

# ---- Slide 20: Demos (light, like other Lattice content slides)
s = slide()
kicker(s, 'Act VI — Live Proof')
title(s, 'Live: 1970s software on the board')
sh = rect(s, 10.9, 0.55, 1.8, 0.55, YELLOW)
shape_text(sh, 'LIVE', size=20, bold=True, color=CHARCOAL)
bullets(s, 0.7, 2.0, 11.9, 4.8, [
    '1.  **Pi** — 50 digits via Machin’s formula, multiprecision arithmetic (2013 program for SCELBI hardware)',
    '2.  **SCELBI Calculator (1974)** — 23-bit floating point:  `12.2 X 2.2 = 26.84`',
    '3.  **HEXPAWN (1973)** — self-modifying code: it learns as it loses',
    '4.  **Finale — the tiny OS:**  power on → BASIC → `MON` → monitor dumps the tokenized program from RAM → `G 1FB6` → back in BASIC, program intact',
], size=18, space=16)
notes(s, 'Rig (rig slide cut — say it here): Versa board, minicom 115200 8N1, local echo '
         'off, DEL for rubout; monitor commands D/W/L/G; DIP switches = interrupts + '
         'READY/WAIT.\n'
         'DEMO CHEAT SHEET (verify commands before talk):\n'
         '• Pi: L in monitor → send hex via send_hex.py → G 0040. ~75 s at 500 kHz for 50 '
         'digits — narrate history while it grinds. 49/50 digits correct; the 50th is the '
         'ORIGINAL program’s guard-byte truncation — great trivia beat.\n'
         '• Calc: load, G, type 12.2 X 2.2 = live.\n'
         '• HEXPAWN: play 2 quick games, point out it stops making the losing move. DROP '
         'CANDIDATE if running long.\n'
         '• Tiny OS: power-cycle board into b8008_basic bitstream. Type a 3-line FOR/NEXT '
         'program, RUN, MON, D the program region, G 1FB6, LIST — program survived.\n'
         '• Fallback: pre-recorded terminal captures of all four in case hardware misbehaves.')

# ---- Slide 22: Takeaways
s = slide()
kicker(s, 'Act VI — Takeaways')
title(s, 'What I’d tell you to steal')
bullets(s, 0.6, 1.95, 12.1, 4.9, [
    '- **Dumb modules, smart control** — structure beats cleverness; the block diagram was right all along',
    '- **Simulation passing is a hypothesis. Hardware is the proof.**',
    '- **AI pair programming works when you own the ground truth** — I specified and verified; it multiplied',
    '- **Old software is the best test suite** — it encodes assumptions no one thinks to write down',
    '- 1972 engineers did all of this with 3,500 transistors. **Respect.**',
], size=19, space=15)
notes(s, 'Close on the human point: the 8008’s designers had no simulator, no regression '
         'suite, no second chances on a mask set.')

# ---- Slide 23: Q&A
s = slide(dark=True)
title(s, 'Questions?', dark=True, size=40, y=1.1)
bullets(s, 0.7, 2.6, 7.6, 3.6, [
    '**Robert Rico**',
    'Project: `github.com/robertrico/intel-8008-vhdl`',
    'SCELBAL: Jim Loos (`jim11662418/8008-SBC`)',
    'Period software: Mike Willegal’s SCELBI archive (willegal.net)',
    'Original SCELBAL: Arnold & Wadsworth, 1976',
], dark=True, size=16, space=10)
tb, tf = box(s, 0.7, 6.3, 8.0, 0.6)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'Board stays up after — come play with BASIC.'
r.font.name = FONT; r.font.size = Pt(19); r.font.bold = True; r.font.italic = True
r.font.color.rgb = YELLOW
qw, qh = fit(REPO + '/docs/images/tiny_os_scelbal.png', 4.3, 3.4)
pic(s, REPO + '/docs/images/tiny_os_scelbal.png', 13.333 - qw - 0.55, 7.5 - qh - 0.5, w=qw)
notes(s, 'Fill in real repo URL if/when public. Invite people to type at the machine — most '
         'memorable 8-bit demo is the one they drive.')

prs.save(OUT)
print('saved', OUT, '-', len(prs.slides._sldIdLst), 'slides')
